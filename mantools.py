#!/usr/bin/env python3
"""
Mantools - merge, convert, and extract PDF & Office files on your own machine.

No network access, no uploads, no telemetry. Everything happens on disk.
The "-> PDF" conversions drive the copy of Microsoft Office already installed
on this Windows machine; nothing is ever sent anywhere.

GUI (light/dark theme, tabbed):
    python3 mantools.py

    Merge tab   - combine PDFs, reorder, per-file rotation, page ranges.
                  Click "Preview" (Ctrl+P) to see every output page as a
                  thumbnail and toggle each one Included/Excluded before
                  merging.
    Convert tab - PDF <-> Word / PowerPoint / Excel, one file or a batch.
    Extract tab - pull pages into a new PDF, or extract text or images.
    Watermark tab - stamp a text or image/logo watermark (centered or tiled)
                  onto a PDF or an image file, with a live preview.

Command line - merge:
    python3 mantools.py -o merged.pdf a.pdf b.pdf
    python3 mantools.py -o merged.pdf "report.pdf:1-3,7" scan.pdf:5-1
    python3 mantools.py --preview "report.pdf:1-3,7" scan.pdf   # dry run

Command line - convert:
    python3 mantools.py convert report.pdf --to docx
    python3 mantools.py convert slides.pdf --to pptx -o out\\slides.pptx
    python3 mantools.py convert a.docx b.xlsx c.pptx --to pdf -o out_dir

Command line - extract:
    python3 mantools.py extract report.pdf --pages 1-3,7 --what pages
    python3 mantools.py extract report.pdf --what text -o notes.txt
    python3 mantools.py extract report.pdf --what images -o out_folder

Command line - watermark (PDF or image):
    python3 mantools.py watermark report.pdf --text CONFIDENTIAL
    python3 mantools.py watermark photo.jpg --text DRAFT --color Red --opacity 20
    python3 mantools.py watermark photo.png --image logo.png --tiled --density 6

Page syntax: 1-3,7 (pages 1,2,3,7)  |  4- (page 4 to end)  |  -2 (start to 2)
             5-1 (pages 5,4,3,2,1 - reverses that span)

Requires: Python 3.8+ and pypdf              ->  pip install pypdf
Optional (unlock more features):
    PyMuPDF   visual previews, PDF -> PPTX/XLSX  ->  pip install pymupdf
    pdf2docx  high-fidelity PDF -> Word          ->  pip install pdf2docx
    python-pptx / openpyxl  PDF -> PPTX / XLSX   ->  pip install python-pptx openpyxl
    comtypes + Microsoft Office (Windows)  Office -> PDF  ->  pip install comtypes
"""

from __future__ import annotations

import os
import re
import sys
import base64
import queue
import tempfile
import threading
from dataclasses import dataclass

APP_NAME = "Mantools"
VERSION = "4.0"
APP_TAGLINE = "Modern PDF utility for merging, converting, and more."

try:
    from pypdf import PdfReader, PdfWriter
    PYPDF_ERROR = None
except ImportError as exc:  # pragma: no cover - environment dependent
    PdfReader = PdfWriter = None  # type: ignore[assignment]
    PYPDF_ERROR = str(exc)

# PyMuPDF is optional. When present, the preview shows real page thumbnails;
# when absent, the preview falls back to an outline with text snippets.
try:
    import pymupdf as _fitz  # modern import name
    HAS_FITZ = True
    # Some dependencies (e.g. pdf2docx) still `import fitz`, whose compat shim
    # prints a deprecation notice. Point that name at pymupdf to avoid it.
    sys.modules.setdefault("fitz", _fitz)
except Exception:
    try:
        import fitz as _fitz  # legacy import name
        HAS_FITZ = True
    except Exception:
        _fitz = None
        HAS_FITZ = False

PYPDF_HINT = (
    "This app needs the pypdf library.\n\n"
    "Install it with:\n"
    "    pip install pypdf\n\n"
    "No internet on this machine? On any machine with a connection run\n"
    "    pip download pypdf -d wheels\n"
    "copy the wheels folder over, then run\n"
    "    pip install --no-index --find-links wheels pypdf"
)

# Optional libraries that power the file converters. Each is probed once so the
# UI can enable only the conversions this machine can actually perform.
IS_WINDOWS = sys.platform.startswith("win")

try:
    import docx as _docx  # python-docx
    HAS_DOCX = True
except Exception:
    _docx = None
    HAS_DOCX = False

try:
    import pptx as _pptx  # python-pptx
    HAS_PPTX = True
except Exception:
    _pptx = None
    HAS_PPTX = False

try:
    import openpyxl as _openpyxl
    HAS_OPENPYXL = True
except Exception:
    _openpyxl = None
    HAS_OPENPYXL = False

try:
    import io as _io
    import contextlib as _contextlib
    # pdf2docx imports the legacy `fitz` shim, which writes a deprecation
    # notice to stderr on first import; swallow just that during import.
    with _contextlib.redirect_stderr(_io.StringIO()):
        import pdf2docx as _pdf2docx  # high-fidelity PDF -> Word
    HAS_PDF2DOCX = True
except Exception:
    _pdf2docx = None
    HAS_PDF2DOCX = False

try:
    import logging as _logging
    # comtypes chats at INFO on import; keep the console quiet.
    _logging.getLogger("comtypes").setLevel(_logging.WARNING)
    import comtypes.client as _comtypes_client  # drives MS Office for -> PDF
    HAS_COMTYPES = True
except Exception:
    _comtypes_client = None
    HAS_COMTYPES = False


@_contextlib.contextmanager
def _quiet_output():
    """Silence chatty third-party INFO logging and direct stdout/stderr writes.

    pdf2docx logs its progress and PyMuPDF prints one-off notices; neither is
    useful to the end user, so wrap those calls in this to keep output clean.
    """
    import logging
    previous = logging.root.manager.disable
    logging.disable(logging.INFO)
    try:
        with _contextlib.redirect_stdout(_io.StringIO()), \
                _contextlib.redirect_stderr(_io.StringIO()):
            yield
    finally:
        logging.disable(previous)

FITZ_HINT = (
    "Install PyMuPDF to see real page thumbnails in the preview:\n"
    "    pip install pymupdf"
)


# --------------------------------------------------------------------------
# Core engine (no GUI - importable and testable on its own)
# --------------------------------------------------------------------------

class MergeError(Exception):
    """Raised with a message that is safe to show the user directly."""


@dataclass
class MergeItem:
    """One source file and how much of it to use."""
    path: str
    pages: str = ""          # "" or "all" means every page
    rotate: int = 0          # 0, 90, 180 or 270, clockwise
    password: str = ""       # for encrypted sources
    page_count: int = 0      # filled in when the file is inspected

    @property
    def name(self) -> str:
        return os.path.basename(self.path)


@dataclass
class PlanPage:
    """One page in the merged output, resolved from a MergeItem."""
    output_index: int        # zero-based position in the merged document
    item: MergeItem          # the source file this page came from
    source_index: int        # zero-based page index within that source
    rotate: int              # clockwise rotation applied to this page
    section_start: bool      # True on the first page taken from each file


def parse_page_range(spec: str, page_count: int) -> list[int]:
    """Turn a page spec into zero-based page indices.

    >>> parse_page_range("1-3,7", 10)
    [0, 1, 2, 6]
    >>> parse_page_range("3-1", 10)
    [2, 1, 0]
    """
    text = (spec or "").strip().lower()
    if not text or text in ("all", "*"):
        return list(range(page_count))

    indices: list[int] = []
    for chunk in text.split(","):
        chunk = chunk.strip()
        if not chunk:
            continue

        if "-" in chunk:
            left, _, right = chunk.partition("-")
            start = _as_page_number(left, default=1, spec=chunk)
            end = _as_page_number(right, default=page_count, spec=chunk)
        else:
            start = end = _as_page_number(chunk, default=None, spec=chunk)

        for value in (start, end):
            if value < 1 or value > page_count:
                raise MergeError(
                    f"Page {value} is out of range - this file has "
                    f"{page_count} page{'s' if page_count != 1 else ''}."
                )

        step = 1 if end >= start else -1
        indices.extend(range(start - 1, end - 1 + step, step))

    if not indices:
        raise MergeError(f"'{spec}' does not select any pages.")
    return indices


def _as_page_number(text: str, default, spec: str) -> int:
    text = text.strip()
    if not text:
        if default is None:
            raise MergeError(f"'{spec}' is not a valid page range.")
        return default
    if not text.isdigit():
        raise MergeError(
            f"'{spec}' is not a valid page range. Use numbers like 1-3,7"
        )
    return int(text)


def count_pages(path: str, password: str = "") -> int:
    """Open a PDF and return its page count. Raises MergeError on failure."""
    _require_pypdf()
    handle = None
    try:
        handle = open(path, "rb")
        reader = PdfReader(handle)
        _unlock(reader, path, password)
        return len(reader.pages)
    except MergeError:
        raise
    except FileNotFoundError:
        raise MergeError(f"{os.path.basename(path)} could not be found.")
    except Exception as exc:
        raise MergeError(f"{os.path.basename(path)} could not be read: {exc}")
    finally:
        if handle is not None:
            handle.close()


def is_encrypted(path: str) -> bool:
    """True if the file needs a password. False if it opens freely."""
    _require_pypdf()
    try:
        with open(path, "rb") as handle:
            reader = PdfReader(handle)
            if not reader.is_encrypted:
                return False
            try:
                return not reader.decrypt("")
            except Exception:
                return True
    except Exception:
        return False


def build_merge_plan(items: list[MergeItem]) -> list[PlanPage]:
    """Resolve `items` into the exact ordered list of pages the merge writes.

    Uses each item's recorded page_count, so it needs no open file handles and
    is cheap enough to run on every edit. Raises MergeError (naming the file)
    if a page range does not fit its source.
    """
    plan: list[PlanPage] = []
    out = 0
    for item in items:
        try:
            indices = parse_page_range(item.pages, item.page_count)
        except MergeError as exc:
            raise MergeError(f"{item.name}: {exc}")
        for position, source_index in enumerate(indices):
            plan.append(PlanPage(
                output_index=out,
                item=item,
                source_index=source_index,
                rotate=item.rotate,
                section_start=(position == 0),
            ))
            out += 1
    return plan


def render_page_png(path: str, page_index: int, rotate: int = 0,
                    password: str = "", target_width: int = 168) -> bytes | None:
    """Render one page to PNG bytes at roughly `target_width` px.

    Returns None if PyMuPDF is unavailable or the page cannot be rendered, so
    callers can fall back gracefully. Never raises.
    """
    if not HAS_FITZ:
        return None
    doc = None
    try:
        doc = _fitz.open(path)
        if doc.needs_pass:
            doc.authenticate(password or "")
        page = doc[page_index]
        if rotate:
            page.set_rotation((page.rotation + rotate) % 360)
        width = max(float(page.rect.width), 1.0)
        zoom = max(min(target_width / width, 4.0), 0.05)
        pixmap = page.get_pixmap(matrix=_fitz.Matrix(zoom, zoom), alpha=False)
        return pixmap.tobytes("png")
    except Exception:
        return None
    finally:
        if doc is not None:
            try:
                doc.close()
            except Exception:
                pass


def extract_snippet(reader, page_index: int, limit: int = 220) -> str:
    """Best-effort single-line text preview of a page. Never raises."""
    try:
        text = reader.pages[page_index].extract_text() or ""
    except Exception:
        return ""
    text = " ".join(text.split())
    if len(text) > limit:
        text = text[:limit].rstrip() + "\u2026"
    return text


def merge_pdfs(
    items: list[MergeItem],
    output_path: str,
    add_bookmarks: bool = True,
    compress: bool = False,
    progress=None,
) -> int:
    """Merge `items` into `output_path`. Returns the page count written.

    Writes to a temporary file first, so a failure part-way through leaves
    the existing output (and every input) untouched. Because of that the
    output may safely be one of the inputs.
    """
    _require_pypdf()
    if not items:
        raise MergeError("Add at least one PDF before merging.")

    handles: list = []
    plan: list[tuple[MergeItem, "PdfReader", list[int]]] = []

    try:
        # Pass 1: open everything and resolve page selections, so an error in
        # the last file is reported before any work is done.
        for item in items:
            try:
                handle = open(item.path, "rb")
            except FileNotFoundError:
                raise MergeError(f"{item.name} could not be found.")
            except OSError as exc:
                raise MergeError(f"{item.name} could not be opened: {exc}")
            handles.append(handle)

            try:
                reader = PdfReader(handle)
                _unlock(reader, item.path, item.password)
                total = len(reader.pages)
            except MergeError:
                raise
            except Exception as exc:
                raise MergeError(f"{item.name} is not a readable PDF: {exc}")

            try:
                indices = parse_page_range(item.pages, total)
            except MergeError as exc:
                raise MergeError(f"{item.name}: {exc}")

            item.page_count = total
            plan.append((item, reader, indices))

        total_pages = sum(len(indices) for _, _, indices in plan)

        # Pass 2: build the merged document.
        writer = PdfWriter()
        done = 0
        for item, reader, indices in plan:
            start_index = done
            for index in indices:
                page = reader.pages[index]
                if item.rotate:
                    page = page.rotate(item.rotate)
                writer.add_page(page)
                if compress:
                    try:
                        writer.pages[-1].compress_content_streams()
                    except Exception:
                        pass  # a page that will not compress is still valid
                done += 1
                if progress:
                    progress(done, total_pages, item.name)
            if add_bookmarks and indices:
                try:
                    writer.add_outline_item(os.path.splitext(item.name)[0],
                                            start_index)
                except Exception:
                    pass  # bookmarks are a nicety, never a reason to fail

        writer.add_metadata({"/Producer": f"{APP_NAME} {VERSION}"})

        directory = os.path.dirname(os.path.abspath(output_path)) or "."
        os.makedirs(directory, exist_ok=True)
        fd, temp_path = tempfile.mkstemp(suffix=".pdf", dir=directory)
        try:
            with os.fdopen(fd, "wb") as out:
                writer.write(out)
        except Exception:
            _quiet_remove(temp_path)
            raise
    finally:
        for handle in handles:
            try:
                handle.close()
            except Exception:
                pass

    try:
        os.replace(temp_path, output_path)
    except OSError as exc:
        _quiet_remove(temp_path)
        raise MergeError(f"Could not save to {output_path}: {exc}")

    return total_pages


def merge_selected_pages(pages, output_path: str, add_bookmarks: bool = True,
                         compress: bool = False, progress=None) -> int:
    """Merge an explicit, ordered list of PlanPage objects into one PDF.

    Used by the preview window, where each page can be individually included
    or excluded, so the selection is no longer whole-file.
    """
    _require_pypdf()
    pages = list(pages)
    if not pages:
        raise MergeError("No pages are selected to merge.")

    readers: dict = {}
    handles: list = []
    try:
        writer = PdfWriter()
        total = len(pages)
        previous_path = None
        for done, plan_page in enumerate(pages, start=1):
            item = plan_page.item
            reader = readers.get(item.path)
            if reader is None:
                try:
                    handle = open(item.path, "rb")
                except OSError as exc:
                    raise MergeError(f"{item.name} could not be opened: {exc}")
                handles.append(handle)
                reader = PdfReader(handle)
                _unlock(reader, item.path, item.password)
                readers[item.path] = reader

            page = reader.pages[plan_page.source_index]
            if plan_page.rotate:
                page = page.rotate(plan_page.rotate)
            writer.add_page(page)
            if compress:
                try:
                    writer.pages[-1].compress_content_streams()
                except Exception:
                    pass
            if add_bookmarks and item.path != previous_path:
                try:
                    writer.add_outline_item(os.path.splitext(item.name)[0],
                                            len(writer.pages) - 1)
                except Exception:
                    pass
            previous_path = item.path
            if progress:
                progress(done, total, item.name)

        writer.add_metadata({"/Producer": f"{APP_NAME} {VERSION}"})
        _write_pdf_atomic(writer, output_path, MergeError)
    finally:
        for handle in handles:
            try:
                handle.close()
            except Exception:
                pass
    return len(pages)


def _unlock(reader, path: str, password: str) -> None:
    if not reader.is_encrypted:
        return
    for candidate in ("", password):
        try:
            if reader.decrypt(candidate):
                return
        except Exception:
            break
    raise MergeError(
        f"{os.path.basename(path)} is password protected. "
        "Add the password and try again."
    )


def _require_pypdf() -> None:
    if PdfReader is None:
        raise MergeError(PYPDF_HINT)


def _quiet_remove(path: str) -> None:
    try:
        os.remove(path)
    except OSError:
        pass


# --------------------------------------------------------------------------
# Conversion engine (PDF <-> Office, all local)
# --------------------------------------------------------------------------

class ConversionError(Exception):
    """Raised with a message that is safe to show the user directly."""


@dataclass
class Route:
    """One supported conversion, e.g. PDF -> Word."""
    key: str                 # stable id, e.g. "pdf2docx"
    label: str               # human label for menus
    src_exts: tuple          # accepted input extensions (lowercase, with dot)
    dst_ext: str             # output extension (with dot)
    func: object             # engine callable (src, out, progress) -> None
    available: bool          # can this machine run it right now?
    hint: str                # why not / what it needs

    def output_for(self, src: str, out_dir: str = "") -> str:
        stem = os.path.splitext(os.path.basename(src))[0]
        folder = out_dir or os.path.dirname(os.path.abspath(src))
        return os.path.join(folder, stem + self.dst_ext)


def _require(condition: bool, message: str) -> None:
    if not condition:
        raise ConversionError(message)


def _com_registered(prog_id: str) -> bool:
    """True if a COM ProgID (e.g. 'Word.Application') is registered."""
    if not IS_WINDOWS:
        return False
    try:
        import winreg
        winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, prog_id).Close()
        return True
    except OSError:
        return False


def _office_ok(prog_id: str) -> bool:
    return IS_WINDOWS and HAS_COMTYPES and _com_registered(prog_id)


# ---- PDF -> Office --------------------------------------------------------

def pdf_to_docx(src: str, out: str, progress=None) -> None:
    """PDF -> Word. Uses pdf2docx (layout-aware) when present, else text-only."""
    if HAS_PDF2DOCX:
        from pdf2docx import Converter
        if progress:
            progress(0, 1, "Analysing layout...")
        converter = Converter(src)
        try:
            with _quiet_output():
                converter.convert(out)
        except Exception as exc:
            raise ConversionError(f"Could not convert to Word: {exc}")
        finally:
            converter.close()
        if progress:
            progress(1, 1, "")
        return

    _require(HAS_FITZ and HAS_DOCX,
             "PDF to Word needs pdf2docx (best) or PyMuPDF + python-docx.\n"
             "    pip install pdf2docx")
    from docx import Document
    document = _fitz.open(src)
    out_doc = Document()
    total = document.page_count
    for index in range(total):
        text = document[index].get_text("text").strip()
        for block in text.split("\n\n"):
            out_doc.add_paragraph(block)
        if index < total - 1:
            out_doc.add_page_break()
        if progress:
            progress(index + 1, total, f"Page {index + 1}")
    out_doc.save(out)


def pdf_to_pptx(src: str, out: str, dpi: int = 150, progress=None) -> None:
    """PDF -> PowerPoint. Each page becomes a full-bleed slide image."""
    _require(HAS_FITZ and HAS_PPTX,
             "PDF to PowerPoint needs PyMuPDF + python-pptx.\n"
             "    pip install pymupdf python-pptx")
    from io import BytesIO
    from pptx import Presentation
    document = _fitz.open(src)
    presentation = Presentation()
    emu_per_point = 12700
    blank = presentation.slide_layouts[6]
    total = document.page_count
    if total == 0:
        raise ConversionError("This PDF has no pages.")
    for index in range(total):
        page = document[index]
        rect = page.rect
        if index == 0:
            presentation.slide_width = int(rect.width * emu_per_point)
            presentation.slide_height = int(rect.height * emu_per_point)
        pixmap = page.get_pixmap(dpi=dpi)
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            BytesIO(pixmap.tobytes("png")), 0, 0,
            width=presentation.slide_width, height=presentation.slide_height)
        if progress:
            progress(index + 1, total, f"Page {index + 1}")
    presentation.save(out)


def pdf_to_xlsx(src: str, out: str, progress=None) -> None:
    """PDF -> Excel. Extracts detected tables; falls back to text per page."""
    _require(HAS_FITZ and HAS_OPENPYXL,
             "PDF to Excel needs PyMuPDF + openpyxl.\n"
             "    pip install pymupdf openpyxl")
    from openpyxl import Workbook
    document = _fitz.open(src)
    workbook = Workbook()
    workbook.remove(workbook.active)
    total = document.page_count
    tables_found = 0
    for index in range(total):
        page = document[index]
        try:
            with _quiet_output():
                tables = list(page.find_tables().tables)
        except Exception:
            tables = []
        for number, table in enumerate(tables, start=1):
            sheet = workbook.create_sheet(title=_sheet_name(index + 1, number))
            for row in table.extract():
                sheet.append(["" if cell is None else str(cell) for cell in row])
            tables_found += 1
        if progress:
            progress(index + 1, total, f"Page {index + 1}")
    if tables_found == 0:
        sheet = workbook.create_sheet(title="Text")
        for index in range(total):
            for line in document[index].get_text("text").splitlines():
                sheet.append([line])
    if not workbook.sheetnames:
        workbook.create_sheet(title="Empty")
    workbook.save(out)


def _sheet_name(page: int, table: int) -> str:
    return f"P{page}_T{table}"[:31]


# ---- Office -> PDF (via installed Microsoft Office) -----------------------

def office_to_pdf(src: str, out: str, progress=None) -> None:
    ext = os.path.splitext(src)[1].lower()
    if progress:
        progress(0, 1, "Opening in Office...")
    if ext in (".doc", ".docx", ".rtf", ".odt"):
        _word_to_pdf(src, out)
    elif ext in (".ppt", ".pptx", ".odp"):
        _powerpoint_to_pdf(src, out)
    elif ext in (".xls", ".xlsx", ".xlsm", ".csv", ".ods"):
        _excel_to_pdf(src, out)
    else:
        raise ConversionError(f"Cannot convert {ext or 'this file'} to PDF.")
    if progress:
        progress(1, 1, "")


def _require_com() -> None:
    _require(IS_WINDOWS and HAS_COMTYPES,
             "Converting to PDF uses Microsoft Office through Windows COM.\n"
             "This needs Windows with Office and the comtypes package.\n"
             "    pip install comtypes")


def _word_to_pdf(src: str, out: str) -> None:
    _require_com()
    _require(_com_registered("Word.Application"),
             "Microsoft Word is not installed, so Word files cannot be "
             "converted to PDF on this machine.")
    import comtypes.client
    word = comtypes.client.CreateObject("Word.Application")
    word.Visible = False
    document = None
    try:
        document = word.Documents.Open(os.path.abspath(src), ReadOnly=True)
        document.SaveAs(os.path.abspath(out), FileFormat=17)  # wdFormatPDF
    except Exception as exc:
        raise ConversionError(f"Word could not create the PDF: {exc}")
    finally:
        if document is not None:
            try:
                document.Close(False)
            except Exception:
                pass
        try:
            word.Quit()
        except Exception:
            pass


def _powerpoint_to_pdf(src: str, out: str) -> None:
    _require_com()
    _require(_com_registered("PowerPoint.Application"),
             "Microsoft PowerPoint is not installed, so slides cannot be "
             "converted to PDF on this machine.")
    import comtypes.client
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    presentation = None
    try:
        presentation = powerpoint.Presentations.Open(
            os.path.abspath(src), WithWindow=False)
        presentation.SaveAs(os.path.abspath(out), 32)  # ppSaveAsPDF
    except Exception as exc:
        raise ConversionError(f"PowerPoint could not create the PDF: {exc}")
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        try:
            powerpoint.Quit()
        except Exception:
            pass


def _excel_to_pdf(src: str, out: str) -> None:
    _require_com()
    _require(_com_registered("Excel.Application"),
             "Microsoft Excel is not installed, so spreadsheets cannot be "
             "converted to PDF on this machine.")
    import comtypes.client
    excel = comtypes.client.CreateObject("Excel.Application")
    excel.Visible = False
    excel.DisplayAlerts = False
    workbook = None
    try:
        workbook = excel.Workbooks.Open(os.path.abspath(src))
        workbook.ExportAsFixedFormat(0, os.path.abspath(out))  # xlTypePDF
    except Exception as exc:
        raise ConversionError(f"Excel could not create the PDF: {exc}")
    finally:
        if workbook is not None:
            try:
                workbook.Close(False)
            except Exception:
                pass
        try:
            excel.Quit()
        except Exception:
            pass


# ---- routes + dispatch ----------------------------------------------------

def conversion_routes() -> list[Route]:
    """Every conversion the app knows about, with live availability."""
    return [
        Route("pdf2docx", "PDF -> Word (.docx)", (".pdf",), ".docx",
              pdf_to_docx, HAS_FITZ and HAS_DOCX,
              "Needs PyMuPDF and python-docx (pdf2docx for best layout)."),
        Route("pdf2pptx", "PDF -> PowerPoint (.pptx)", (".pdf",), ".pptx",
              pdf_to_pptx, HAS_FITZ and HAS_PPTX,
              "Needs PyMuPDF and python-pptx."),
        Route("pdf2xlsx", "PDF -> Excel (.xlsx)", (".pdf",), ".xlsx",
              pdf_to_xlsx, HAS_FITZ and HAS_OPENPYXL,
              "Needs PyMuPDF and openpyxl."),
        Route("docx2pdf", "Word -> PDF", (".doc", ".docx"), ".pdf",
              office_to_pdf, _office_ok("Word.Application"),
              "Needs Microsoft Word (Windows)."),
        Route("pptx2pdf", "PowerPoint -> PDF", (".ppt", ".pptx"), ".pdf",
              office_to_pdf, _office_ok("PowerPoint.Application"),
              "Needs Microsoft PowerPoint (Windows)."),
        Route("xlsx2pdf", "Excel -> PDF", (".xls", ".xlsx", ".xlsm", ".csv"),
              ".pdf", office_to_pdf, _office_ok("Excel.Application"),
              "Needs Microsoft Excel (Windows)."),
    ]


def route_by_key(key: str) -> "Route | None":
    for route in conversion_routes():
        if route.key == key:
            return route
    return None


def pick_route(src: str, target: str) -> "Route | None":
    """Choose the conversion for a source file and a target format word."""
    ext = os.path.splitext(src)[1].lower()
    dst = "." + target.lower().lstrip(".")
    for route in conversion_routes():
        if route.dst_ext == dst and ext in route.src_exts:
            return route
    return None


def convert_file(src: str, route: Route, out_dir: str = "",
                 out_path: str = "", progress=None) -> str:
    """Run one conversion. Returns the output path actually written."""
    if not os.path.isfile(src):
        raise ConversionError(f"{os.path.basename(src)} could not be found.")
    if not route.available:
        raise ConversionError(route.hint)
    target = out_path or route.output_for(src, out_dir)
    directory = os.path.dirname(os.path.abspath(target)) or "."
    os.makedirs(directory, exist_ok=True)
    route.func(src, target, progress=progress)
    return target


# --------------------------------------------------------------------------
# Extract engine (pull pages, text, or images out of a PDF)
# --------------------------------------------------------------------------

class ExtractError(Exception):
    """Raised with a message that is safe to show the user directly."""


def _resolve_indices(pages: str, total: int) -> list[int]:
    try:
        return parse_page_range(pages, total)
    except MergeError as exc:
        raise ExtractError(str(exc))


def _write_pdf_atomic(writer, output_path: str, error_cls=MergeError) -> None:
    directory = os.path.dirname(os.path.abspath(output_path)) or "."
    os.makedirs(directory, exist_ok=True)
    fd, temp_path = tempfile.mkstemp(suffix=".pdf", dir=directory)
    try:
        with os.fdopen(fd, "wb") as out:
            writer.write(out)
    except Exception:
        _quiet_remove(temp_path)
        raise
    try:
        os.replace(temp_path, output_path)
    except OSError as exc:
        _quiet_remove(temp_path)
        raise error_cls(f"Could not save to {output_path}: {exc}")


def extract_pages(src: str, out: str, pages: str = "all",
                  password: str = "", progress=None) -> int:
    """Extract selected pages into a new PDF. Returns the page count written."""
    _require_pypdf()
    if not os.path.isfile(src):
        raise ExtractError(f"{os.path.basename(src)} could not be found.")
    handle = open(src, "rb")
    try:
        try:
            reader = PdfReader(handle)
            _unlock(reader, src, password)
            total = len(reader.pages)
        except MergeError as exc:
            raise ExtractError(str(exc))
        except Exception as exc:
            raise ExtractError(f"{os.path.basename(src)} is not a readable PDF: {exc}")
        indices = _resolve_indices(pages, total)
        writer = PdfWriter()
        for done, index in enumerate(indices, start=1):
            writer.add_page(reader.pages[index])
            if progress:
                progress(done, len(indices), f"Page {index + 1}")
        writer.add_metadata({"/Producer": f"{APP_NAME} {VERSION}"})
        _write_pdf_atomic(writer, out, ExtractError)
    finally:
        handle.close()
    return len(indices)


def extract_text(src: str, out: str, pages: str = "all",
                 password: str = "", progress=None) -> int:
    """Extract page text into a UTF-8 .txt file. Returns characters written."""
    if not os.path.isfile(src):
        raise ExtractError(f"{os.path.basename(src)} could not be found.")
    chunks: list[str] = []
    if HAS_FITZ:
        document = _fitz.open(src)
        if document.needs_pass:
            document.authenticate(password or "")
        total = document.page_count
        indices = _resolve_indices(pages, total)
        for done, index in enumerate(indices, start=1):
            chunks.append(document[index].get_text("text"))
            if progress:
                progress(done, len(indices), f"Page {index + 1}")
    else:
        _require_pypdf()
        with open(src, "rb") as handle:
            reader = PdfReader(handle)
            _unlock(reader, src, password)
            total = len(reader.pages)
            indices = _resolve_indices(pages, total)
            for done, index in enumerate(indices, start=1):
                chunks.append(reader.pages[index].extract_text() or "")
                if progress:
                    progress(done, len(indices), f"Page {index + 1}")
    text = "\n\n".join(chunks)
    directory = os.path.dirname(os.path.abspath(out)) or "."
    os.makedirs(directory, exist_ok=True)
    with open(out, "w", encoding="utf-8") as fh:
        fh.write(text)
    return len(text)


def extract_images(src: str, out_dir: str, pages: str = "all",
                   password: str = "", progress=None) -> int:
    """Save embedded images from the chosen pages. Returns how many were saved."""
    _require(HAS_FITZ, "Extracting images needs PyMuPDF.\n    pip install pymupdf")
    if not os.path.isfile(src):
        raise ExtractError(f"{os.path.basename(src)} could not be found.")
    document = _fitz.open(src)
    if document.needs_pass:
        document.authenticate(password or "")
    total = document.page_count
    indices = _resolve_indices(pages, total)
    os.makedirs(out_dir, exist_ok=True)
    stem = os.path.splitext(os.path.basename(src))[0]
    seen: set = set()
    saved = 0
    for done, index in enumerate(indices, start=1):
        for image in document[index].get_images(full=True):
            xref = image[0]
            if xref in seen:
                continue
            seen.add(xref)
            try:
                info = document.extract_image(xref)
            except Exception:
                continue
            data = info.get("image")
            if not data:
                continue
            saved += 1
            ext = info.get("ext", "png")
            name = f"{stem}_p{index + 1}_{saved:03d}.{ext}"
            with open(os.path.join(out_dir, name), "wb") as fh:
                fh.write(data)
        if progress:
            progress(done, len(indices), f"Page {index + 1}")
    return saved


# --------------------------------------------------------------------------
# Watermark engine (stamp text across PDF pages)
# --------------------------------------------------------------------------

class WatermarkError(Exception):
    """Raised with a message that is safe to show the user directly."""


WATERMARK_COLORS = {
    "Gray":  (0.50, 0.50, 0.50),
    "Red":   (0.80, 0.16, 0.16),
    "Blue":  (0.16, 0.35, 0.82),
    "Green": (0.13, 0.53, 0.27),
    "Black": (0.00, 0.00, 0.00),
}


def add_watermark(src: str, out: str, text: str = "", pages: str = "all",
                  font_size: float = 48, color=(0.5, 0.5, 0.5),
                  opacity: float = 0.15, angle: float = 45,
                  tiled: bool = False, image: str = "", image_width: float = 220,
                  density: int = 4, password: str = "", progress=None) -> int:
    """Stamp a text or image watermark across the chosen pages.

    Pass `image` (a PNG/JPG path) for a logo watermark; otherwise `text` is
    stamped. `density` sets how many copies fit across the page when `tiled`.
    Returns the number of pages stamped.
    """
    _require(HAS_FITZ, "Watermarking needs PyMuPDF.\n    pip install pymupdf")
    use_image = bool(image)
    if use_image:
        if not os.path.isfile(image):
            raise WatermarkError(
                f"Watermark image {os.path.basename(image)} could not be found.")
        try:
            png_bytes, (img_w, img_h) = _prepare_wm_image(image, opacity, angle)
        except Exception as exc:
            raise WatermarkError(f"Could not read the watermark image: {exc}")
    elif not (text or "").strip():
        raise WatermarkError("Enter some watermark text.")
    if not os.path.isfile(src):
        raise WatermarkError(f"{os.path.basename(src)} could not be found.")
    try:
        document = _fitz.open(src)
    except Exception as exc:
        raise WatermarkError(f"{os.path.basename(src)} is not a readable PDF: {exc}")
    if document.needs_pass:
        document.authenticate(password or "")
    total = document.page_count
    try:
        wanted = set(_resolve_indices(pages, total))
    except ExtractError as exc:
        document.close()
        raise WatermarkError(str(exc))

    done = 0
    for index in range(total):
        if index not in wanted:
            continue
        if use_image:
            _stamp_image_page(document[index], png_bytes, img_w, img_h,
                              float(image_width), tiled, density)
        else:
            _stamp_page(document[index], text, float(font_size), color,
                        float(opacity), float(angle), tiled, density)
        done += 1
        if progress:
            progress(done, len(wanted), f"Page {index + 1}")

    directory = os.path.dirname(os.path.abspath(out)) or "."
    os.makedirs(directory, exist_ok=True)
    fd, tmp = tempfile.mkstemp(suffix=".pdf", dir=directory)
    os.close(fd)
    try:
        document.save(tmp, garbage=3, deflate=True)
    finally:
        document.close()
    try:
        os.replace(tmp, out)
    except OSError as exc:
        _quiet_remove(tmp)
        raise WatermarkError(f"Could not save to {out}: {exc}")
    return done


def _prepare_wm_image(path: str, opacity: float, angle: float):
    """Load a watermark image, bake in opacity + rotation. Returns (png, size)."""
    from PIL import Image
    import io
    image = Image.open(path).convert("RGBA")
    if opacity < 1:
        alpha = image.getchannel("A").point(lambda v: int(v * opacity))
        image.putalpha(alpha)
    turn = float(angle) % 360
    if turn:
        image = image.rotate(turn, expand=True, resample=Image.BICUBIC)
    buffer = io.BytesIO()
    image.save(buffer, "PNG")
    return buffer.getvalue(), image.size


def _stamp_image_page(page, png_bytes, img_w, img_h, target_w, tiled,
                      density=4) -> None:
    rect = page.rect
    aspect = (img_h / img_w) if img_w else 1.0
    box_w = target_w
    box_h = target_w * aspect
    if tiled:
        cols = max(1, int(round(density)))
        gap = rect.width / cols
        y = gap * 0.5
        while y < rect.height + gap * 0.5:
            x = gap * 0.5
            while x < rect.width + gap * 0.5:
                page.insert_image(_fitz.Rect(x - box_w / 2, y - box_h / 2,
                                             x + box_w / 2, y + box_h / 2),
                                  stream=png_bytes, keep_proportion=True,
                                  overlay=True)
                x += gap
            y += gap
    else:
        cx, cy = rect.width / 2, rect.height / 2
        page.insert_image(_fitz.Rect(cx - box_w / 2, cy - box_h / 2,
                                     cx + box_w / 2, cy + box_h / 2),
                          stream=png_bytes, keep_proportion=True, overlay=True)


def _stamp_page(page, text, font_size, color, opacity, angle, tiled,
                density=4) -> None:
    rect = page.rect
    fontname = "helv"
    matrix = _fitz.Matrix(1, 1).prerotate(angle)
    text_w = _fitz.get_text_length(text, fontname=fontname, fontsize=font_size)
    kw = dict(fontname=fontname, fontsize=font_size, color=color,
              fill_opacity=opacity, overlay=True)
    if tiled:
        cols = max(1, int(round(density)))
        gap = rect.width / cols
        y = gap * 0.5
        while y < rect.height + gap * 0.5:
            x = gap * 0.5
            while x < rect.width + gap * 0.5:
                point = _fitz.Point(x - text_w / 2, y + font_size * 0.30)
                page.insert_text(point, text, morph=(_fitz.Point(x, y), matrix),
                                 **kw)
                x += gap
            y += gap
    else:
        cx, cy = rect.width / 2, rect.height / 2
        point = _fitz.Point(cx - text_w / 2, cy + font_size * 0.35)
        page.insert_text(point, text, morph=(_fitz.Point(cx, cy), matrix), **kw)


WATERMARK_IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif",
                        ".tiff", ".webp")


def is_watermark_image(path: str) -> bool:
    """True if `path` looks like an image file (vs. a PDF) to watermark."""
    return os.path.splitext(path)[1].lower() in WATERMARK_IMAGE_EXTS


def _image_to_page(src: str, ref_width: float = 612.0):
    """Wrap an image in a one-page PDF sized to a reference width.

    Watermark sizes then behave the same as on a PDF page regardless of the
    image's pixel dimensions. Returns (doc, page, scale) where rendering the
    page at `scale` restores the original resolution.
    """
    try:
        pix = _fitz.Pixmap(src)
    except Exception as exc:
        raise WatermarkError(
            f"{os.path.basename(src)} is not a readable image: {exc}")
    img_w, img_h = pix.width, pix.height
    if not img_w or not img_h:
        raise WatermarkError(f"{os.path.basename(src)} has no pixels.")
    scale = img_w / ref_width
    page_w, page_h = ref_width, img_h / scale
    doc = _fitz.open()
    page = doc.new_page(width=page_w, height=page_h)
    page.insert_image(_fitz.Rect(0, 0, page_w, page_h), pixmap=pix)
    return doc, page, scale


def watermark_image(src: str, out: str, text: str = "", pages: str = "all",
                    font_size: float = 48, color=(0.5, 0.5, 0.5),
                    opacity: float = 0.15, angle: float = 45,
                    tiled: bool = False, image: str = "",
                    image_width: float = 220, density: int = 4,
                    password: str = "", progress=None) -> int:
    """Stamp a text or image watermark onto an image file. Returns 1.

    `pages`/`password` are accepted (for a common signature with
    `add_watermark`) but ignored - an image is a single frame.
    """
    _require(HAS_FITZ, "Watermarking needs PyMuPDF.\n    pip install pymupdf")
    use_image = bool(image)
    if use_image:
        if not os.path.isfile(image):
            raise WatermarkError(
                f"Watermark image {os.path.basename(image)} could not be found.")
    elif not (text or "").strip():
        raise WatermarkError("Enter some watermark text.")
    if not os.path.isfile(src):
        raise WatermarkError(f"{os.path.basename(src)} could not be found.")

    doc, page, scale = _image_to_page(src)
    try:
        if use_image:
            png_bytes, (img_w, img_h) = _prepare_wm_image(image, opacity, angle)
            _stamp_image_page(page, png_bytes, img_w, img_h, float(image_width),
                              tiled, density)
        else:
            _stamp_page(page, text, float(font_size), color, float(opacity),
                        float(angle), tiled, density)
        result = page.get_pixmap(matrix=_fitz.Matrix(scale, scale), alpha=False)
    finally:
        doc.close()

    directory = os.path.dirname(os.path.abspath(out)) or "."
    os.makedirs(directory, exist_ok=True)
    from PIL import Image
    pil = Image.frombytes("RGB", (result.width, result.height), result.samples)
    ext = os.path.splitext(out)[1].lower()
    try:
        if ext in (".jpg", ".jpeg"):
            pil.save(out, quality=92)
        else:
            pil.save(out)
    except Exception as exc:
        raise WatermarkError(f"Could not save to {out}: {exc}")
    if progress:
        progress(1, 1, os.path.basename(src))
    return 1


def apply_watermark(src: str, out: str, **kwargs) -> int:
    """Watermark a PDF or an image, chosen by the source file type."""
    if is_watermark_image(src):
        return watermark_image(src, out, **kwargs)
    return add_watermark(src, out, **kwargs)


# --------------------------------------------------------------------------
# Command line
# --------------------------------------------------------------------------

RANGE_SUFFIX = re.compile(r"^[0-9]+(\s*-\s*[0-9]*)?(\s*,\s*[0-9]+(\s*-\s*[0-9]*)?)*$")


def split_path_and_range(argument: str) -> tuple[str, str]:
    """Split 'file.pdf:1-3' into ('file.pdf', '1-3').

    Splits on the last colon and only when the tail looks like a page range,
    so Windows paths such as C:\\docs\\a.pdf survive intact.
    """
    head, sep, tail = argument.rpartition(":")
    if sep and tail and (RANGE_SUFFIX.match(tail.strip()) or tail.strip() == "all"):
        return head, tail.strip()
    return argument, ""


def run_cli(argv: list[str]) -> int:
    import argparse

    parser = argparse.ArgumentParser(
        prog="mantools.py",
        description=f"{APP_NAME} - merge PDFs locally. Run with no arguments "
                    "for the window version.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="Page syntax:  file.pdf:1-3,7   file.pdf:4-   file.pdf:5-1 (reversed)",
    )
    parser.add_argument("inputs", nargs="+", metavar="FILE[:PAGES]")
    parser.add_argument("-o", "--output", metavar="FILE",
                        help="where to write the merged PDF")
    parser.add_argument("--preview", action="store_true",
                        help="show the resolved page plan and exit (writes nothing)")
    parser.add_argument("--rotate", type=int, default=0, choices=[0, 90, 180, 270],
                        help="rotate every page clockwise by this many degrees")
    parser.add_argument("--password", default="",
                        help="password for encrypted inputs")
    parser.add_argument("--no-bookmarks", action="store_true",
                        help="skip the per-file bookmarks in the merged file")
    parser.add_argument("--compress", action="store_true",
                        help="compress page content to shrink the output")
    parser.add_argument("--quiet", action="store_true", help="print nothing on success")
    args = parser.parse_args(argv)

    items = []
    for raw in args.inputs:
        path, pages = split_path_and_range(raw)
        items.append(MergeItem(path=path, pages=pages, rotate=args.rotate,
                               password=args.password))

    if args.preview:
        return _cli_preview(items)

    if not args.output:
        print("Error: -o/--output is required (or use --preview for a dry run).",
              file=sys.stderr)
        return 2

    def show(done, total, name):
        if not args.quiet:
            sys.stdout.write(f"\r  {done}/{total} pages - {name[:40]}")
            sys.stdout.flush()

    try:
        pages = merge_pdfs(
            items,
            args.output,
            add_bookmarks=not args.no_bookmarks,
            compress=args.compress,
            progress=show,
        )
    except MergeError as exc:
        if not args.quiet:
            sys.stdout.write("\r")
        print(f"Error: {exc}", file=sys.stderr)
        return 1

    if not args.quiet:
        size = os.path.getsize(args.output) / 1_048_576
        print(f"\rMerged {len(items)} file{'s' if len(items) != 1 else ''} "
              f"into {args.output} - {pages} pages, {size:.1f} MB")
    return 0


def _cli_preview(items: list[MergeItem]) -> int:
    """Print the resolved merge plan without writing anything."""
    for item in items:
        try:
            item.page_count = count_pages(item.path, item.password)
        except MergeError as exc:
            print(f"Error: {exc}", file=sys.stderr)
            return 1
    try:
        plan = build_merge_plan(items)
    except MergeError as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1

    print(f"Preview - {len(plan)} page{'s' if len(plan) != 1 else ''} from "
          f"{len(items)} file{'s' if len(items) != 1 else ''} (nothing written)\n")
    print(f"  {'out':>4}  {'source':<34} {'src page':>8}  rotate")
    print(f"  {'-'*4}  {'-'*34} {'-'*8}  {'-'*6}")
    for page in plan:
        marker = "> " if page.section_start else "  "
        rot = f"{page.rotate} deg" if page.rotate else ""
        name = page.item.name
        if len(name) > 34:
            name = name[:31] + "..."
        print(f"  {page.output_index + 1:>4}  {marker}{name:<32} "
              f"{page.source_index + 1:>8}  {rot:>6}")
    print("\n  '>' marks where a new file (and bookmark) begins.")
    return 0


def run_convert_cli(argv: list[str]) -> int:
    import argparse

    parser = argparse.ArgumentParser(
        prog="mantools.py convert",
        description=f"{APP_NAME} - convert between PDF and Office formats, "
                    "entirely on this machine.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="Examples:\n"
               "  mantools.py convert report.pdf --to docx\n"
               "  mantools.py convert slides.pdf --to pptx -o out\\slides.pptx\n"
               "  mantools.py convert a.docx b.xlsx c.pptx --to pdf -o out_dir",
    )
    parser.add_argument("inputs", nargs="+", metavar="FILE")
    parser.add_argument("--to", required=True, metavar="FORMAT",
                        choices=["pdf", "docx", "pptx", "xlsx"],
                        help="target format: pdf, docx, pptx or xlsx")
    parser.add_argument("-o", "--output", metavar="PATH",
                        help="output file (single input) or folder (many inputs)")
    parser.add_argument("--quiet", action="store_true",
                        help="print nothing on success")
    args = parser.parse_args(argv)

    single = len(args.inputs) == 1
    out_dir, out_file = "", ""
    if args.output:
        looks_like_dir = (os.path.isdir(args.output)
                          or args.output.endswith(("\\", "/"))
                          or not os.path.splitext(args.output)[1])
        if single and not looks_like_dir:
            out_file = args.output
        else:
            out_dir = args.output

    exit_code = 0
    for src in args.inputs:
        route = pick_route(src, args.to)
        base = os.path.basename(src)
        if route is None:
            ext = os.path.splitext(src)[1] or "(no extension)"
            print(f"Error: no converter from {ext} to .{args.to} for {base}.",
                  file=sys.stderr)
            exit_code = 1
            continue
        if not route.available:
            print(f"Error: {route.label} is unavailable. {route.hint}",
                  file=sys.stderr)
            exit_code = 1
            continue

        def show(done, total, name, _label=route.label):
            if not args.quiet:
                tail = f" - {name}" if name else ""
                sys.stdout.write(f"\r  {_label}: {done}/{total}{tail}      ")
                sys.stdout.flush()

        try:
            written = convert_file(src, route, out_dir=out_dir,
                                   out_path=(out_file if single else ""),
                                   progress=show)
        except ConversionError as exc:
            if not args.quiet:
                sys.stdout.write("\r")
            print(f"Error: {exc}", file=sys.stderr)
            exit_code = 1
            continue
        except Exception as exc:  # never lose an unexpected failure
            if not args.quiet:
                sys.stdout.write("\r")
            print(f"Error: unexpected problem with {base}: {exc}", file=sys.stderr)
            exit_code = 1
            continue

        if not args.quiet:
            size = os.path.getsize(written) / 1_048_576
            print(f"\r{base} -> {written}  ({size:.1f} MB)".ljust(70))
    return exit_code


def run_extract_cli(argv: list[str]) -> int:
    import argparse

    parser = argparse.ArgumentParser(
        prog="mantools.py extract",
        description=f"{APP_NAME} - pull pages, text, or images out of a PDF.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="Examples:\n"
               "  mantools.py extract report.pdf --pages 1-3,7 --what pages\n"
               "  mantools.py extract report.pdf --what text -o notes.txt\n"
               "  mantools.py extract report.pdf --what images -o out_folder",
    )
    parser.add_argument("input", metavar="FILE")
    parser.add_argument("--pages", default="all", metavar="RANGE",
                        help="which pages, e.g. 1-3,7 or all (default: all)")
    parser.add_argument("--what", choices=["pages", "text", "images"],
                        default="pages", help="what to extract (default: pages)")
    parser.add_argument("-o", "--output", metavar="PATH",
                        help="output file (pages/text) or folder (images)")
    parser.add_argument("--password", default="")
    parser.add_argument("--quiet", action="store_true")
    args = parser.parse_args(argv)

    src = args.input
    stem = os.path.splitext(os.path.basename(src))[0]
    folder = os.path.dirname(os.path.abspath(src))

    def show(done, total, name):
        if not args.quiet:
            sys.stdout.write(f"\r  {name} {done}/{total}   ")
            sys.stdout.flush()

    try:
        if args.what == "pages":
            out = args.output or os.path.join(folder, f"{stem}_extracted.pdf")
            count = extract_pages(src, out, args.pages, args.password, show)
            result = f"{count} pages -> {out}"
        elif args.what == "text":
            out = args.output or os.path.join(folder, f"{stem}.txt")
            chars = extract_text(src, out, args.pages, args.password, show)
            result = f"{chars} characters -> {out}"
        else:
            out = args.output or os.path.join(folder, f"{stem}_images")
            saved = extract_images(src, out, args.pages, args.password, show)
            result = f"{saved} images -> {out}"
    except (ExtractError, MergeError) as exc:
        if not args.quiet:
            sys.stdout.write("\r")
        print(f"Error: {exc}", file=sys.stderr)
        return 1

    if not args.quiet:
        print(f"\rExtracted {result}".ljust(70))
    return 0


def run_watermark_cli(argv: list[str]) -> int:
    import argparse

    parser = argparse.ArgumentParser(
        prog="mantools.py watermark",
        description=f"{APP_NAME} - stamp a text or image watermark onto a PDF "
                    "or an image file.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="Examples:\n"
               "  mantools.py watermark report.pdf --text CONFIDENTIAL\n"
               "  mantools.py watermark photo.jpg --text DRAFT --color Red "
               "--opacity 20\n"
               "  mantools.py watermark photo.png --image logo.png --tiled",
    )
    parser.add_argument("input", metavar="FILE", help="a PDF or an image file")
    parser.add_argument("--text", default="", help="the watermark text")
    parser.add_argument("--image", default="", metavar="PNG",
                        help="stamp this image/logo instead of text")
    parser.add_argument("--pages", default="all", metavar="RANGE",
                        help="pages to stamp, e.g. 1-3,7 or all (default: all)")
    parser.add_argument("--size", type=float, default=48,
                        help="font size (text) — for images use --image-width")
    parser.add_argument("--image-width", type=float, default=220,
                        help="image watermark width in points (default: 220)")
    parser.add_argument("--color", default="Gray", choices=list(WATERMARK_COLORS),
                        help="text colour (default: Gray)")
    parser.add_argument("--opacity", type=float, default=15,
                        help="opacity percent, 1-100 (default: 15)")
    parser.add_argument("--angle", type=float, default=45,
                        help="rotation in degrees (default: 45)")
    parser.add_argument("--tiled", action="store_true",
                        help="repeat the watermark across the whole page")
    parser.add_argument("--density", type=int, default=4,
                        help="tiled: how many copies across the page (default: 4)")
    parser.add_argument("-o", "--output", metavar="PATH")
    parser.add_argument("--password", default="")
    parser.add_argument("--quiet", action="store_true")
    args = parser.parse_args(argv)

    if not args.text and not args.image:
        print("Error: pass --text or --image.", file=sys.stderr)
        return 2

    src = args.input
    stem = os.path.splitext(os.path.basename(src))[0]
    src_ext = os.path.splitext(src)[1] or ".pdf"
    out = args.output or os.path.join(
        os.path.dirname(os.path.abspath(src)), f"{stem}_watermarked{src_ext}")

    def show(done, total, name):
        if not args.quiet:
            sys.stdout.write(f"\r  {name} {done}/{total}   ")
            sys.stdout.flush()

    try:
        stamped = apply_watermark(
            src, out, text=args.text, pages=args.pages, font_size=args.size,
            color=WATERMARK_COLORS[args.color],
            opacity=max(0.0, min(1.0, args.opacity / 100.0)),
            angle=args.angle, tiled=args.tiled, image=args.image,
            image_width=args.image_width, density=max(1, args.density),
            password=args.password, progress=show)
    except (WatermarkError, MergeError) as exc:
        if not args.quiet:
            sys.stdout.write("\r")
        print(f"Error: {exc}", file=sys.stderr)
        return 1

    if not args.quiet:
        if is_watermark_image(src):
            print(f"\rWatermarked {os.path.basename(src)} -> {out}".ljust(70))
        else:
            print(f"\rWatermarked {stamped} page{'s' if stamped != 1 else ''} -> "
                  f"{out}".ljust(70))
    return 0


# --------------------------------------------------------------------------
# Window version  ("Mantools" - light/dark themed desktop UI)
# --------------------------------------------------------------------------

# Two colour themes; the window switches between them at runtime.
THEMES = {
    "light": {
        "bg": "#eef1f5", "card": "#ffffff", "border": "#dbe0e6",
        "ink": "#111827", "muted": "#667085",
        "accent": "#2f6fed", "accent_dk": "#2258cc", "accent_soft": "#aac3f5",
        "sel": "#e8f0ff", "head": "#f3f5f8", "page": "#ffffff",
        "field": "#ffffff", "field_bd": "#cbd2da",
        "on": "#16a34a", "on_soft": "#e9f7ef", "off": "#c2c8d0",
        "drop": "#f7f9fc", "drop_bd": "#b9c2cf",
        "footer": "#eaf1ff", "toolbtn": "#ffffff",
    },
    "dark": {
        "bg": "#0f141b", "card": "#1a212c", "border": "#2b3543",
        "ink": "#e7ebf1", "muted": "#93a1b3",
        "accent": "#4f8bf7", "accent_dk": "#3b74e0", "accent_soft": "#33456a",
        "sel": "#1f2d46", "head": "#161d27", "page": "#0d131b",
        "field": "#0d131b", "field_bd": "#39465a",
        "on": "#22c55e", "on_soft": "#122a1e", "off": "#3a4658",
        "drop": "#141c26", "drop_bd": "#39465a",
        "footer": "#132444", "toolbtn": "#1f2836",
    },
}

PREVIEW_THUMB_WIDTH = 168
PREVIEW_COLUMNS = 4
PREVIEW_MAX_CARDS = 400  # keep the preview responsive for huge merges


def run_gui(_smoke=None) -> int:
    """Launch the Mantools window (light/dark, tabbed: Merge / Convert / Extract).

    `_smoke` is a testing seam: when given a callable, the app is built and
    handed to it (instead of entering the blocking main loop), then torn down.
    """
    try:
        import tkinter as tk
        from tkinter import (ttk, filedialog, messagebox, simpledialog,
                             colorchooser, font as tkfont)
    except ImportError:
        print(
            "The window version needs tkinter, which is missing from this "
            "Python.\n"
            "  Debian/Ubuntu:  sudo apt install python3-tk\n"
            "  Fedora:         sudo dnf install python3-tkinter\n"
            "  macOS/Windows:  install Python from python.org\n\n"
            "The command line version works without it:\n"
            "  python3 mantools.py -o merged.pdf a.pdf b.pdf",
            file=sys.stderr,
        )
        return 1

    try:
        from tkinterdnd2 import TkinterDnD, DND_FILES
        base_class, drag_and_drop = TkinterDnD.Tk, True
    except Exception:
        base_class, drag_and_drop = tk.Tk, False
        DND_FILES = None

    UI = "Segoe UI"

    def ellipsize(text, limit):
        return text if len(text) <= limit else text[:limit - 1] + "…"

    # ------------------------------------------------------------------
    # A small canvas pill toggle (used per page in the preview).
    # ------------------------------------------------------------------
    class ToggleSwitch(tk.Canvas):
        def __init__(self, master, colors, value=True, command=None):
            super().__init__(master, width=46, height=24, highlightthickness=0,
                             bd=0, bg=colors["card"], cursor="hand2")
            self._c = colors
            self._value = value
            self._command = command
            self.bind("<Button-1>", self._on_click)
            self._draw()

        def _draw(self):
            self.delete("all")
            colour = self._c["on"] if self._value else self._c["off"]
            self.create_oval(2, 2, 22, 22, fill=colour, outline=colour)
            self.create_oval(24, 2, 44, 22, fill=colour, outline=colour)
            self.create_rectangle(12, 2, 34, 22, fill=colour, outline=colour)
            knob = 34 if self._value else 12
            self.create_oval(knob - 10, 3, knob + 10, 21,
                             fill="#ffffff", outline="#ffffff")

        def _on_click(self, _event=None):
            self._value = not self._value
            self._draw()
            if self._command:
                self._command(self._value)

        def get(self):
            return self._value

    # ------------------------------------------------------------------
    # Preview modal - every output page with an Included/Excluded toggle.
    # ------------------------------------------------------------------
    class PreviewWindow(tk.Toplevel):
        def __init__(self, master, plan, colors, on_merge):
            super().__init__(master)
            self.C = colors
            self.plan = plan
            self.on_merge = on_merge
            self.included = [True] * len(plan)
            self.title(f"{APP_NAME} – Preview")
            self.configure(bg=colors["bg"])
            self.minsize(720, 520)
            self.geometry("900x680")

            self._events = queue.Queue()
            self._stop = threading.Event()
            self._images = []
            self._holders = {}
            self._state_labels = {}
            self._poll_id = None

            self._build()
            self.protocol("WM_DELETE_WINDOW", self._close)
            self.bind("<Escape>", lambda _e: self._close())
            self.transient(master)
            self._poll_id = self.after(60, self._drain)
            threading.Thread(target=self._load, daemon=True).start()

        def _build(self):
            C = self.C
            files = len({p.item.path for p in self.plan})
            head = tk.Frame(self, bg=C["bg"])
            head.pack(fill="x", padx=20, pady=(18, 8))
            line = tk.Frame(head, bg=C["bg"])
            line.pack(anchor="w")
            tk.Label(line, text=APP_NAME, bg=C["bg"], fg=C["accent"],
                     font=(UI, 15, "bold")).pack(side="left")
            tk.Label(line, text="  Preview before merge", bg=C["bg"], fg=C["ink"],
                     font=(UI, 15, "bold")).pack(side="left")
            tk.Label(head, bg=C["bg"], fg=C["muted"], font=(UI, 10),
                     text=f"{len(self.plan)} page{'s' if len(self.plan) != 1 else ''} "
                          f"from {files} file{'s' if files != 1 else ''}, shown in the "
                          "exact order they will be written.").pack(anchor="w",
                                                                    pady=(2, 0))

            body = tk.Frame(self, bg=C["bg"])
            body.pack(fill="both", expand=True, padx=12)
            self.canvas = tk.Canvas(body, bg=C["bg"], highlightthickness=0, bd=0)
            vbar = ttk.Scrollbar(body, orient="vertical", command=self.canvas.yview)
            self.canvas.configure(yscrollcommand=vbar.set)
            vbar.pack(side="right", fill="y")
            self.canvas.pack(side="left", fill="both", expand=True)
            self.grid = tk.Frame(self.canvas, bg=C["bg"])
            self._win = self.canvas.create_window((0, 0), window=self.grid,
                                                  anchor="nw")
            self.grid.bind("<Configure>", lambda _e: self.canvas.configure(
                scrollregion=self.canvas.bbox("all")))
            self.canvas.bind("<Configure>", lambda e: self.canvas.itemconfigure(
                self._win, width=e.width))
            self.canvas.bind("<MouseWheel>", self._wheel)
            self._build_cards()

            foot = tk.Frame(self, bg=C["footer"])
            foot.pack(fill="x", side="bottom")
            inner = tk.Frame(foot, bg=C["footer"])
            inner.pack(fill="x", padx=20, pady=12)
            self.count_var = tk.StringVar(value="")
            tk.Label(inner, textvariable=self.count_var, bg=C["footer"], fg=C["ink"],
                     font=(UI, 11, "bold")).pack(side="left")
            ttk.Button(inner, text="Cancel", style="Ghost.TButton",
                       command=self._close).pack(side="right")
            ttk.Button(inner, text="Merge Pages", style="Primary.TButton",
                       command=self._merge).pack(side="right", padx=(0, 8))
            self._update_count()

        def _build_cards(self):
            C = self.C
            for column in range(PREVIEW_COLUMNS):
                self.grid.columnconfigure(column, weight=1, uniform="pv")
            shown = self.plan[:PREVIEW_MAX_CARDS]
            for position, page in enumerate(shown):
                row, column = divmod(position, PREVIEW_COLUMNS)
                card = tk.Frame(self.grid, bg=C["card"], highlightthickness=1,
                                highlightbackground=C["on"], bd=0)
                card.grid(row=row, column=column, padx=8, pady=8, sticky="nsew")
                top = tk.Frame(card, bg=C["card"])
                top.pack(fill="x", padx=10, pady=(8, 4))
                tk.Label(top, text=f"Page {page.output_index + 1}", bg=C["card"],
                         fg=C["ink"], font=(UI, 12, "bold")).pack(side="left")
                switch_wrap = tk.Frame(top, bg=C["card"])
                switch_wrap.pack(side="right")
                state = tk.Label(switch_wrap, text="Included", bg=C["card"],
                                 fg=C["on"], font=(UI, 8, "bold"))
                state.pack(side="left", padx=(0, 6))
                toggle = ToggleSwitch(
                    switch_wrap, C, value=True,
                    command=lambda v, i=page.output_index: self._on_toggle(i, v))
                toggle.pack(side="left")
                self._state_labels[page.output_index] = (state, card)

                holder = tk.Frame(card, bg=C["page"], width=PREVIEW_THUMB_WIDTH,
                                  height=int(PREVIEW_THUMB_WIDTH * 1.3),
                                  highlightthickness=1, highlightbackground=C["border"])
                holder.pack(padx=10, pady=4)
                holder.pack_propagate(False)
                image_label = tk.Label(holder, bg=C["page"], fg=C["muted"],
                                       text="…", font=(UI, 9))
                image_label.pack(fill="both", expand=True)
                self._holders[page.output_index] = image_label

                tk.Label(card, text=ellipsize(page.item.name, 24), bg=C["card"],
                         fg=C["muted"], font=(UI, 8)).pack(fill="x", padx=10,
                                                           pady=(2, 8))
            if len(self.plan) > PREVIEW_MAX_CARDS:
                tk.Label(self.grid, bg=C["bg"], fg=C["muted"], font=(UI, 9),
                         text=f"Showing the first {PREVIEW_MAX_CARDS} of "
                              f"{len(self.plan)} pages.").grid(
                    row=(len(shown) // PREVIEW_COLUMNS) + 1, column=0,
                    columnspan=PREVIEW_COLUMNS, sticky="w", pady=6)

        def _wheel(self, event):
            self.canvas.yview_scroll(-1 if event.delta > 0 else 1, "units")

        def _on_toggle(self, index, value):
            self.included[index] = value
            state, card = self._state_labels[index]
            state.configure(text="Included" if value else "Excluded",
                            fg=self.C["on"] if value else self.C["muted"])
            card.configure(highlightbackground=self.C["on"] if value
                           else self.C["border"])
            holder = self._holders.get(index)
            if holder is not None:
                photo = getattr(holder, "_photo", None)
                if value and photo is not None:
                    holder.configure(image=photo, text="")
                else:
                    holder.configure(image="", text="✕ excluded",
                                     fg=self.C["muted"], font=(UI, 11, "bold"))
            self._update_count()

        def _update_count(self):
            selected = sum(1 for flag in self.included if flag)
            self.count_var.set(f"Total pages selected: {selected}")

        def _load(self):
            shown = self.plan[:PREVIEW_MAX_CARDS]
            for page in shown:
                if self._stop.is_set():
                    break
                png = render_page_png(page.item.path, page.source_index,
                                      page.rotate, page.item.password,
                                      PREVIEW_THUMB_WIDTH)
                self._events.put((page.output_index, png))
            self._events.put(("done", None))

        def _drain(self):
            try:
                while True:
                    index, png = self._events.get_nowait()
                    if index == "done":
                        continue
                    self._set_thumb(index, png)
            except queue.Empty:
                pass
            if self._poll_id is not None:
                self._poll_id = self.after(60, self._drain)

        def _set_thumb(self, index, png):
            holder = self._holders.get(index)
            if holder is None or not png:
                if holder is not None:
                    holder.configure(text="no preview")
                return
            try:
                photo = tk.PhotoImage(
                    master=self, data=base64.b64encode(png).decode("ascii"))
            except Exception:
                holder.configure(text="no preview")
                return
            self._images.append(photo)
            holder._photo = photo
            if self.included[index]:
                holder.configure(image=photo, text="")

        def _merge(self):
            chosen = [self.plan[i] for i, flag in enumerate(self.included) if flag]
            if not chosen:
                messagebox.showinfo("Nothing selected",
                                    "Include at least one page to merge.")
                return
            self._close()
            self.on_merge(chosen)

        def _close(self):
            self._stop.set()
            if self._poll_id is not None:
                try:
                    self.after_cancel(self._poll_id)
                except Exception:
                    pass
                self._poll_id = None
            self.destroy()

    # ------------------------------------------------------------------
    # Main window
    # ------------------------------------------------------------------
    class MantoolsApp(base_class):
        def __init__(self):
            super().__init__()
            self.title(APP_NAME)
            self._set_app_icon()
            self.minsize(900, 720)
            self.geometry("1040x820")
            self.theme = "light"
            self.C = dict(THEMES[self.theme])

            self.items: list[MergeItem] = []
            self.events: queue.Queue = queue.Queue()
            self.working = False
            self._drain_id = None
            self._selected = None
            self._drag_from = None

            self.conv_inputs: list[str] = []
            self.conv_events: queue.Queue = queue.Queue()
            self.conv_working = False
            self._conv_drain_id = None

            self.ex_events: queue.Queue = queue.Queue()
            self.ex_working = False
            self._ex_drain_id = None

            self.wm_events: queue.Queue = queue.Queue()
            self.wm_working = False
            self._wm_drain_id = None

            self.tab = "merge"
            self.routes = conversion_routes()

            self.style = ttk.Style(self)
            if "clam" in self.style.theme_names():
                self.style.theme_use("clam")
            self._init_styles()
            self.configure(bg=self.C["bg"])
            self._build_chrome()
            self._select_tab("merge")
            self._bind_keys()
            self._drain_id = self.after(100, self._drain_events)
            self._conv_drain_id = self.after(100, self._drain_conv_events)
            self._ex_drain_id = self.after(100, self._drain_ex_events)
            self._wm_drain_id = self.after(100, self._drain_wm_events)

            if PYPDF_ERROR:
                self.after(200, self._warn_missing_pypdf)

        # -- app icon --------------------------------------------------
        def _set_app_icon(self):
            """Use the Mantools logo files next to this script, if present."""
            try:
                base = os.path.dirname(os.path.abspath(__file__))
            except Exception:
                return
            png = os.path.join(base, "mantools.png")
            ico = os.path.join(base, "mantools.ico")
            # On Windows the multi-resolution .ico has crisp 16/32px frames;
            # use it for this window and for child Toplevels (e.g. Preview).
            if sys.platform.startswith("win") and os.path.exists(ico):
                try:
                    self.iconbitmap(ico)
                    self.iconbitmap(default=ico)
                    return
                except Exception:
                    pass
            # Elsewhere (or if the .ico is missing) fall back to the PNG.
            if os.path.exists(png):
                try:
                    self._app_icon = tk.PhotoImage(file=png)
                    self.iconphoto(True, self._app_icon)
                except Exception:
                    pass

        # -- theming ---------------------------------------------------
        def _init_styles(self):
            C = self.C
            s = self.style
            s.configure("App.TFrame", background=C["bg"])
            s.configure("Card.TFrame", background=C["card"])
            s.configure("Head.TFrame", background=C["head"])
            s.configure("Sep.TFrame", background=C["border"])
            s.configure("Accent.TFrame", background=C["accent"])
            s.configure("Footer.TFrame", background=C["footer"])

            s.configure("Title.TLabel", background=C["bg"], foreground=C["ink"],
                        font=(UI, 22, "bold"))
            s.configure("Tag.TLabel", background=C["bg"], foreground=C["muted"],
                        font=(UI, 10))
            s.configure("H2.TLabel", background=C["bg"], foreground=C["ink"],
                        font=(UI, 12, "bold"))
            s.configure("Muted.TLabel", background=C["bg"], foreground=C["muted"],
                        font=(UI, 9))
            s.configure("Ink.TLabel", background=C["bg"], foreground=C["ink"],
                        font=(UI, 10))
            s.configure("Field.TLabel", background=C["bg"], foreground=C["ink"],
                        font=(UI, 10))
            s.configure("Hint.TLabel", background=C["bg"], foreground=C["accent"],
                        font=(UI, 9))
            s.configure("Status.TLabel", background=C["bg"], foreground=C["muted"],
                        font=(UI, 10))
            s.configure("CardInk.TLabel", background=C["card"], foreground=C["ink"],
                        font=(UI, 10))
            s.configure("CardMuted.TLabel", background=C["card"],
                        foreground=C["muted"], font=(UI, 9))

            s.configure("Tab.TLabel", background=C["bg"], foreground=C["muted"],
                        font=(UI, 13))
            s.configure("TabActive.TLabel", background=C["bg"],
                        foreground=C["accent"], font=(UI, 13, "bold"))

            s.configure("RowNum.TLabel", background=C["card"], foreground=C["muted"],
                        font=(UI, 10))
            s.configure("RowName.TLabel", background=C["card"], foreground=C["ink"],
                        font=(UI, 10))
            s.configure("RowNameSel.TLabel", background=C["card"],
                        foreground=C["accent"], font=(UI, 10, "bold"))
            s.configure("RowPages.TLabel", background=C["card"],
                        foreground=C["muted"], font=(UI, 10))
            s.configure("Handle.TLabel", background=C["card"], foreground=C["muted"],
                        font=(UI, 13))
            s.configure("RotOff.TLabel", background=C["card"], foreground=C["muted"],
                        font=(UI, 9))
            s.configure("RotOn.TLabel", background=C["card"], foreground=C["accent"],
                        font=(UI, 9, "bold"))
            s.configure("HeadCell.TLabel", background=C["head"],
                        foreground=C["muted"], font=(UI, 9, "bold"))

            for name, bg in (("Primary.TButton", C["accent"]),
                             ("PrimaryBig.TButton", C["accent"])):
                s.configure(name, background=bg, foreground="#ffffff",
                            borderwidth=0, focusthickness=0,
                            font=(UI, 10 if name == "Primary.TButton" else 13,
                                  "bold"),
                            padding=(16, 8) if name == "Primary.TButton"
                            else (20, 14))
                s.map(name, background=[("disabled", C["accent_soft"]),
                                        ("pressed", C["accent_dk"]),
                                        ("active", C["accent_dk"])],
                      foreground=[("disabled", "#eef2ff")])
            s.configure("Ghost.TButton", background=C["card"], foreground=C["ink"],
                        borderwidth=1, focusthickness=0, relief="flat",
                        font=(UI, 10), padding=(12, 7), bordercolor=C["border"])
            s.map("Ghost.TButton", background=[("active", C["head"])],
                  bordercolor=[("active", C["accent"])])
            s.configure("Tool.TButton", background=C["toolbtn"], foreground=C["ink"],
                        borderwidth=1, focusthickness=0, relief="flat",
                        font=(UI, 10), padding=(12, 7), bordercolor=C["border"])
            s.map("Tool.TButton", background=[("active", C["head"])],
                  bordercolor=[("active", C["accent"])])
            s.configure("Seg.TButton", background=C["card"], foreground=C["ink"],
                        borderwidth=1, focusthickness=0, relief="flat",
                        font=(UI, 10), padding=(14, 8), bordercolor=C["border"])
            s.map("Seg.TButton", background=[("active", C["head"])])
            s.configure("SegOn.TButton", background=C["accent"], foreground="#ffffff",
                        borderwidth=1, focusthickness=0, relief="flat",
                        font=(UI, 10, "bold"), padding=(14, 8),
                        bordercolor=C["accent"])
            s.map("SegOn.TButton", background=[("active", C["accent_dk"])])
            s.configure("PillOn.TButton", background=C["accent"], foreground="#ffffff",
                        borderwidth=0, focusthickness=0, font=(UI, 11), padding=(8, 4))
            s.map("PillOn.TButton", background=[("active", C["accent"])])
            s.configure("PillOff.TButton", background=C["card"], foreground=C["muted"],
                        borderwidth=0, focusthickness=0, font=(UI, 11), padding=(8, 4))
            s.map("PillOff.TButton", background=[("active", C["head"])])

            s.configure("Accent.Horizontal.TProgressbar", background=C["accent"],
                        troughcolor=C["head"], borderwidth=0, thickness=8)
            s.configure("TCheckbutton", background=C["bg"], foreground=C["ink"],
                        font=(UI, 10))
            s.map("TCheckbutton", background=[("active", C["bg"])],
                  foreground=[("disabled", C["muted"])])
            s.configure("TCombobox", fieldbackground=C["field"],
                        background=C["field"], foreground=C["ink"],
                        arrowcolor=C["ink"], bordercolor=C["field_bd"], padding=5)
            s.configure("TEntry", fieldbackground=C["field"], foreground=C["ink"],
                        bordercolor=C["field_bd"], insertcolor=C["ink"], padding=6)
            s.configure("TSpinbox", fieldbackground=C["field"],
                        foreground=C["ink"], background=C["field"],
                        arrowcolor=C["ink"], bordercolor=C["field_bd"], padding=4)

        def _apply_theme(self):
            self.C = dict(THEMES[self.theme])
            self._init_styles()
            self.configure(bg=self.C["bg"])
            self._refresh_theme_pill()
            self._update_tab_styles()
            self._select_tab(self.tab)

        def _toggle_theme(self, target):
            if target != self.theme:
                self.theme = target
                self._apply_theme()

        # -- chrome: header + tabs + content ---------------------------
        def _build_chrome(self):
            self.root_frame = ttk.Frame(self, style="App.TFrame",
                                        padding=(24, 18, 24, 12))
            self.root_frame.pack(fill="both", expand=True)

            header = ttk.Frame(self.root_frame, style="App.TFrame")
            header.pack(fill="x")
            titles = ttk.Frame(header, style="App.TFrame")
            titles.pack(side="left")
            ttk.Label(titles, text=APP_NAME, style="Title.TLabel").pack(anchor="w")
            ttk.Label(titles, text=APP_TAGLINE, style="Tag.TLabel").pack(
                anchor="w", pady=(2, 0))

            self.pill = ttk.Frame(header, style="Card.TFrame", padding=2)
            self.pill.pack(side="right", anchor="n")
            self.btn_light = ttk.Button(self.pill, text="☀", width=3,
                                        command=lambda: self._toggle_theme("light"))
            self.btn_light.pack(side="left")
            self.btn_dark = ttk.Button(self.pill, text="☾", width=3,
                                       command=lambda: self._toggle_theme("dark"))
            self.btn_dark.pack(side="left")
            self._refresh_theme_pill()

            tabbar = ttk.Frame(self.root_frame, style="App.TFrame")
            tabbar.pack(fill="x", pady=(16, 0))
            self._tabs = {}
            for key, text in (("merge", "Merge"), ("convert", "Convert"),
                              ("extract", "Extract"), ("watermark", "Watermark")):
                holder = ttk.Frame(tabbar, style="App.TFrame")
                holder.pack(side="left", padx=(0, 26))
                label = ttk.Label(holder, text=text, style="Tab.TLabel",
                                  cursor="hand2")
                label.pack()
                underline = ttk.Frame(holder, style="Accent.TFrame", height=3)
                underline.pack(fill="x", pady=(6, 0))
                label.bind("<Button-1>", lambda _e, k=key: self._select_tab(k))
                self._tabs[key] = (label, underline)
            ttk.Frame(self.root_frame, style="Sep.TFrame", height=1).pack(
                fill="x", pady=(0, 0))

            self.content = ttk.Frame(self.root_frame, style="App.TFrame",
                                     padding=(0, 14, 0, 0))
            self.content.pack(fill="both", expand=True)

        def _refresh_theme_pill(self):
            self.pill.configure(style="Card.TFrame")
            self.btn_light.configure(
                style="PillOn.TButton" if self.theme == "light" else "PillOff.TButton")
            self.btn_dark.configure(
                style="PillOn.TButton" if self.theme == "dark" else "PillOff.TButton")

        def _update_tab_styles(self):
            for key, (label, underline) in self._tabs.items():
                active = (key == self.tab)
                label.configure(style="TabActive.TLabel" if active else "Tab.TLabel")
                underline.configure(style="Accent.TFrame" if active else "App.TFrame")

        def _select_tab(self, key):
            self.tab = key
            self._update_tab_styles()
            for child in self.content.winfo_children():
                child.destroy()
            if key == "merge":
                self._build_merge(self.content)
            elif key == "convert":
                self._build_convert(self.content)
            elif key == "extract":
                self._build_extract(self.content)
            else:
                self._build_watermark(self.content)

        def _bind_keys(self):
            self.bind("<Control-o>", lambda _e: self._route_add())
            self.bind("<Control-p>", lambda _e: self.tab == "merge"
                      and self.open_preview())

        def _route_add(self):
            if self.tab == "merge":
                self.add_files()
            elif self.tab == "convert":
                self.conv_add_files()

        def _warn_missing_pypdf(self):
            messagebox.showwarning("pypdf is not installed", PYPDF_HINT)

        def _bind_wheel(self, widget, canvas):
            widget.bind("<MouseWheel>",
                        lambda e: canvas.yview_scroll(-1 if e.delta > 0 else 1,
                                                      "units"))

        # ==============================================================
        # MERGE TAB
        # ==============================================================
        def _build_merge(self, parent):
            C = self.C
            bar = ttk.Frame(parent, style="App.TFrame")
            bar.pack(fill="x")
            ttk.Label(bar, text="", style="Muted.TLabel").pack(side="left")
            ttk.Button(bar, text="↓ Down", style="Tool.TButton",
                       command=lambda: self.move(1)).pack(side="right")
            ttk.Button(bar, text="↑ Up", style="Tool.TButton",
                       command=lambda: self.move(-1)).pack(side="right", padx=(0, 8))
            ttk.Button(bar, text="\U0001f5d1 Clear", style="Tool.TButton",
                       command=self.clear_all).pack(side="right", padx=(0, 8))
            ttk.Button(bar, text="⊕ Add Files", style="Tool.TButton",
                       command=self.add_files).pack(side="right", padx=(0, 8))

            card = ttk.Frame(parent, style="Card.TFrame", padding=1)
            card.pack(fill="both", expand=True, pady=(12, 0))

            header = ttk.Frame(card, style="Head.TFrame", padding=(14, 8))
            header.pack(fill="x")
            self._table_columns(header, "HeadCell.TLabel",
                                 ("Order", "File Name", "Pages", "Rotate", ""))

            body = ttk.Frame(card, style="Card.TFrame")
            body.pack(fill="both", expand=True)
            self.mcanvas = tk.Canvas(body, bg=C["card"], highlightthickness=0, bd=0)
            mscroll = ttk.Scrollbar(body, orient="vertical",
                                    command=self.mcanvas.yview)
            self.mcanvas.configure(yscrollcommand=mscroll.set)
            mscroll.pack(side="right", fill="y")
            self.mcanvas.pack(side="left", fill="both", expand=True)
            self.mrows = ttk.Frame(self.mcanvas, style="Card.TFrame")
            self._mwin = self.mcanvas.create_window((0, 0), window=self.mrows,
                                                    anchor="nw")
            self.mrows.bind("<Configure>", lambda _e: self.mcanvas.configure(
                scrollregion=self.mcanvas.bbox("all")))
            self.mcanvas.bind("<Configure>", lambda e: self.mcanvas.itemconfigure(
                self._mwin, width=e.width))
            self._bind_wheel(self.mcanvas, self.mcanvas)

            controls = ttk.Frame(parent, style="App.TFrame")
            controls.pack(fill="x", pady=(14, 0))
            controls.columnconfigure(0, weight=1)
            controls.columnconfigure(1, weight=1)

            left = ttk.Frame(controls, style="App.TFrame")
            left.grid(row=0, column=0, sticky="nsew", padx=(0, 16))
            ttk.Label(left, text="Page Range:", style="Field.TLabel").grid(
                row=0, column=0, sticky="w", pady=(0, 6))
            self.pages_var = tk.StringVar()
            page_entry = ttk.Entry(left, textvariable=self.pages_var)
            page_entry.grid(row=0, column=1, sticky="ew", padx=(10, 0), pady=(0, 6))
            page_entry.bind("<Return>", lambda _e: self.apply_pages())
            page_entry.bind("<FocusOut>", lambda _e: self.apply_pages())
            ttk.Label(left, text="Default Rotation:", style="Field.TLabel").grid(
                row=1, column=0, sticky="w", pady=(0, 8))
            self.default_rot = tk.StringVar(value="0° (None)")
            rot = ttk.Combobox(left, textvariable=self.default_rot, state="readonly",
                               width=14, values=["0° (None)", "90°",
                                                 "180°", "270°"])
            rot.grid(row=1, column=1, sticky="w", padx=(10, 0), pady=(0, 8))
            rot.bind("<<ComboboxSelected>>", lambda _e: self.apply_default_rotation())
            left.columnconfigure(1, weight=1)
            self.bookmarks_var = tk.BooleanVar(value=False)
            self.compress_var = tk.BooleanVar(value=False)
            ttk.Checkbutton(left, text="Bookmark each file in the result",
                            variable=self.bookmarks_var).grid(
                row=2, column=0, columnspan=2, sticky="w")
            ttk.Checkbutton(left, text="Compress output",
                            variable=self.compress_var).grid(
                row=3, column=0, columnspan=2, sticky="w", pady=(2, 0))

            right = ttk.Frame(controls, style="App.TFrame")
            right.grid(row=0, column=1, sticky="nsew")
            right.columnconfigure(0, weight=1)
            self.merge_button = ttk.Button(right, text="Merge PDF",
                                           style="PrimaryBig.TButton",
                                           command=self.start_merge)
            self.merge_button.grid(row=0, column=0, columnspan=2, sticky="ew")
            self.progress = ttk.Progressbar(
                right, mode="determinate",
                style="Accent.Horizontal.TProgressbar")
            self.progress.grid(row=1, column=0, columnspan=2, sticky="ew",
                               pady=(8, 6))
            self.ready_var = tk.StringVar(value="")
            ttk.Label(right, textvariable=self.ready_var, style="Ink.TLabel").grid(
                row=2, column=0, sticky="w")
            self.preview_button = ttk.Button(right, text="Preview",
                                              style="Ghost.TButton",
                                              command=self.open_preview)
            self.preview_button.grid(row=2, column=1, sticky="e")

            ttk.Frame(parent, style="Sep.TFrame", height=1).pack(fill="x",
                                                                 pady=(14, 0))
            self.status_var = tk.StringVar(value="No files yet.")
            ttk.Label(parent, textvariable=self.status_var, style="Status.TLabel"
                      ).pack(anchor="w", pady=(8, 0))

            if drag_and_drop:
                self.mcanvas.drop_target_register(DND_FILES)
                self.mcanvas.dnd_bind("<<Drop>>", self._on_drop)
            self.merge_refresh()

        def _table_columns(self, parent, style, texts):
            widths = (70, 0, 100, 260, 34)
            ttk.Frame(parent, style="Head.TFrame", width=3).grid(row=0, column=0)
            for col, (text, width) in enumerate(zip(texts, widths), start=1):
                parent.columnconfigure(col, weight=(1 if width == 0 else 0),
                                       minsize=width)
                ttk.Label(parent, text=text, style=style,
                          background=self.C["head"]).grid(
                    row=0, column=col, sticky="w", padx=(0, 8))

        def add_files(self):
            paths = filedialog.askopenfilenames(
                title="Add PDFs",
                filetypes=[("PDF files", "*.pdf"), ("All files", "*.*")])
            self._add_paths(paths)

        def _on_drop(self, event):
            self._add_paths(self.mcanvas.tk.splitlist(event.data))

        def _add_paths(self, paths):
            added, skipped = 0, []
            for path in paths:
                path = os.path.abspath(path)
                if not path.lower().endswith(".pdf") or not os.path.isfile(path):
                    skipped.append(os.path.basename(path))
                    continue
                password = ""
                if PdfReader is not None and is_encrypted(path):
                    password = simpledialog.askstring(
                        "Password needed",
                        f"{os.path.basename(path)} is protected.\nEnter its password:",
                        show="•", parent=self) or ""
                item = MergeItem(path=path, password=password)
                try:
                    item.page_count = count_pages(path, password)
                except MergeError as exc:
                    skipped.append(f"{os.path.basename(path)} ({exc})")
                    continue
                self.items.append(item)
                added += 1
            self.merge_refresh()
            if skipped:
                messagebox.showwarning(
                    "Some files were not added",
                    "These could not be used:\n\n" + "\n".join(skipped[:10]))
            elif added:
                self.status_var.set(f"Added {added} file"
                                    f"{'s' if added != 1 else ''}.")

        def remove_selected(self):
            if self._selected is None:
                return
            del self.items[self._selected]
            self._selected = None
            self.merge_refresh()

        def clear_all(self):
            self.items.clear()
            self._selected = None
            self.merge_refresh()

        def move(self, step):
            if self._selected is None:
                return
            target = self._selected + step
            if 0 <= target < len(self.items):
                self.items[self._selected], self.items[target] = (
                    self.items[target], self.items[self._selected])
                self._selected = target
                self.merge_refresh()

        def select_row(self, index):
            self._selected = index
            self.pages_var.set(self.items[index].pages)
            self.merge_refresh()

        def set_rotate(self, index, degrees):
            self.items[index].rotate = degrees
            self.merge_refresh()

        def apply_pages(self):
            if self._selected is None:
                return
            spec = self.pages_var.get().strip()
            item = self.items[self._selected]
            try:
                parse_page_range(spec, item.page_count)
            except MergeError as exc:
                messagebox.showerror("That page range will not work",
                                     f"{item.name}: {exc}")
                return
            item.pages = spec
            self.merge_refresh()

        def apply_default_rotation(self):
            raw = self.default_rot.get()
            degrees = int("".join(ch for ch in raw if ch.isdigit()) or "0")
            for item in self.items:
                item.rotate = degrees
            self.merge_refresh()

        def _drag_start(self, index):
            self._drag_from = index

        def _drag_drop(self, event):
            if self._drag_from is None:
                return
            y = self.mcanvas.canvasy(
                event.y_root - self.mcanvas.winfo_rooty())
            target = min(len(self.items) - 1, max(0, int(y // 44)))
            if target != self._drag_from:
                item = self.items.pop(self._drag_from)
                self.items.insert(target, item)
                self._selected = target
            self._drag_from = None
            self.merge_refresh()

        def merge_refresh(self):
            if not hasattr(self, "mrows"):
                return
            for child in self.mrows.winfo_children():
                child.destroy()
            widths = (70, 0, 100, 260, 34)
            total = 0
            for index, item in enumerate(self.items):
                try:
                    selected_pages = len(parse_page_range(item.pages,
                                                          item.page_count))
                except MergeError:
                    selected_pages = 0
                total += selected_pages
                row = ttk.Frame(self.mrows, style="Card.TFrame", padding=(0, 6))
                row.grid(row=index, column=0, sticky="ew")
                self.mrows.columnconfigure(0, weight=1)
                inner = ttk.Frame(row, style="Card.TFrame")
                inner.pack(fill="x", padx=14)
                chosen = (index == self._selected)
                stripe = ttk.Frame(inner,
                                   style="Accent.TFrame" if chosen else "Card.TFrame",
                                   width=3)
                stripe.grid(row=0, column=0, sticky="ns", padx=(0, 8))
                for col, width in enumerate(widths, start=1):
                    inner.columnconfigure(col, weight=(1 if width == 0 else 0),
                                          minsize=width)
                ttk.Label(inner, text=f"{index + 1}.", style="RowNum.TLabel").grid(
                    row=0, column=1, sticky="w")
                name = ttk.Label(inner, text=ellipsize(item.name, 46),
                                 style="RowNameSel.TLabel" if chosen
                                 else "RowName.TLabel", cursor="hand2")
                name.grid(row=0, column=2, sticky="w")
                name.bind("<Button-1>", lambda _e, i=index: self.select_row(i))
                pages_label = (f"({item.page_count} page"
                               f"{'s' if item.page_count != 1 else ''})")
                ttk.Label(inner, text=pages_label, style="RowPages.TLabel").grid(
                    row=0, column=3, sticky="w")
                rotate = ttk.Frame(inner, style="Card.TFrame")
                rotate.grid(row=0, column=4, sticky="w")
                for degrees in (0, 90, 180, 270):
                    lab = ttk.Label(rotate, text=f"↻{degrees}°",
                                    style="RotOn.TLabel" if item.rotate == degrees
                                    else "RotOff.TLabel", cursor="hand2")
                    lab.pack(side="left", padx=(0, 10))
                    lab.bind("<Button-1>",
                             lambda _e, i=index, d=degrees: self.set_rotate(i, d))
                handle = ttk.Label(inner, text="≡", style="Handle.TLabel",
                                   cursor="fleur")
                handle.grid(row=0, column=5, sticky="e")
                handle.bind("<ButtonPress-1>", lambda _e, i=index: self._drag_start(i))
                handle.bind("<ButtonRelease-1>", self._drag_drop)
                sep = ttk.Frame(self.mrows, style="Sep.TFrame", height=1)
                sep.grid(row=index, column=0, sticky="ews")

            if self._selected is not None and self._selected >= len(self.items):
                self._selected = None
            if self.items:
                self.status_var.set(
                    f"Added {len(self.items)} file"
                    f"{'s' if len(self.items) != 1 else ''}. Total pages: {total}.")
                self.ready_var.set(f"Ready to merge {len(self.items)} file"
                                   f"{'s' if len(self.items) != 1 else ''}.")
                self.merge_button.state(["!disabled"])
            else:
                self.status_var.set("No files yet.")
                self.ready_var.set("Add PDFs to begin.")
                self.merge_button.state(["disabled"])
                self.pages_var.set("")

        # -- preview + merge -------------------------------------------
        def open_preview(self):
            if self.working or not self.items:
                if not self.items:
                    messagebox.showinfo("Nothing to preview",
                                        "Add at least one PDF first.")
                return
            try:
                plan = build_merge_plan(self.items)
            except MergeError as exc:
                messagebox.showerror("That page range will not work", str(exc))
                return
            if not plan:
                messagebox.showinfo("Nothing to preview",
                                    "The current page ranges select no pages.")
                return
            PreviewWindow(self, plan, self.C, on_merge=self.merge_pages)

        def _ask_output(self):
            path = filedialog.asksaveasfilename(
                title="Save merged PDF as", defaultextension=".pdf",
                initialfile="merged.pdf", filetypes=[("PDF files", "*.pdf")])
            return path

        def start_merge(self):
            if self.working:
                return
            if not self.items:
                messagebox.showinfo("Nothing to merge", "Add at least one PDF first.")
                return
            output = self._ask_output()
            if not output:
                return
            items = [MergeItem(it.path, it.pages, it.rotate, it.password,
                               it.page_count) for it in self.items]
            bookmarks = self.bookmarks_var.get()
            compress = self.compress_var.get()
            self._launch_merge(lambda progress: merge_pdfs(
                items, output, add_bookmarks=bookmarks,
                compress=compress, progress=progress), output)

        def merge_pages(self, pages):
            if self.working:
                return
            output = self._ask_output()
            if not output:
                return
            bookmarks = self.bookmarks_var.get()
            compress = self.compress_var.get()
            self._launch_merge(lambda progress: merge_selected_pages(
                pages, output, add_bookmarks=bookmarks,
                compress=compress, progress=progress), output)

        def _launch_merge(self, job, output):
            self.working = True
            self.merge_button.state(["disabled"])
            self.progress.configure(value=0, maximum=100)
            self.status_var.set("Merging…")

            def work():
                try:
                    pages = job(lambda done, total, name:
                                self.events.put(("progress", done, total, name)))
                    self.events.put(("done", pages, output, None))
                except MergeError as exc:
                    self.events.put(("done", 0, output, str(exc)))
                except Exception as exc:
                    self.events.put(("done", 0, output, f"Unexpected problem: {exc}"))

            threading.Thread(target=work, daemon=True).start()

        def _drain_events(self):
            try:
                while True:
                    event = self.events.get_nowait()
                    if event[0] == "progress":
                        _, done, total, name = event
                        self.progress.configure(maximum=max(total, 1), value=done)
                        self.status_var.set(f"Merging {name} · {done}/{total} pages")
                    else:
                        _, pages, output, error = event
                        self.working = False
                        self.merge_button.state(["!disabled"])
                        if error:
                            self.progress.configure(value=0)
                            self.status_var.set("Merge failed.")
                            messagebox.showerror("Merge failed", error)
                        else:
                            size = os.path.getsize(output) / 1_048_576
                            self.progress.configure(value=self.progress["maximum"])
                            self.status_var.set(
                                f"Saved {os.path.basename(output)} · "
                                f"{pages} pages · {size:.1f} MB")
                            if messagebox.askyesno(
                                    "Merged", f"Saved {pages} pages to:\n{output}\n\n"
                                    "Open the containing folder?"):
                                open_folder(output)
            except queue.Empty:
                pass
            self._drain_id = self.after(100, self._drain_events)

        # ==============================================================
        # CONVERT TAB
        # ==============================================================
        def _conv_route(self):
            label = self.conv_route_var.get()
            for route in self.routes:
                if route.label == label:
                    return route
            return self.routes[0]

        def _build_convert(self, parent):
            C = self.C
            top = ttk.Frame(parent, style="App.TFrame")
            top.pack(fill="x")
            ttk.Label(top, text="Conversion Type:", style="Field.TLabel").pack(
                side="left", padx=(0, 10))
            self.conv_route_var = tk.StringVar(
                value=getattr(self, "conv_route_var", None) and
                self.conv_route_var.get() or self.routes[0].label)
            combo = ttk.Combobox(top, textvariable=self.conv_route_var,
                                 state="readonly", width=34,
                                 values=[r.label for r in self.routes])
            combo.pack(side="left")
            combo.bind("<<ComboboxSelected>>", lambda _e: self._conv_route_changed())

            self.conv_hint = tk.StringVar()
            ttk.Label(parent, textvariable=self.conv_hint, style="Hint.TLabel").pack(
                anchor="w", pady=(8, 0))

            self.drop_wrap = ttk.Frame(parent, style="Card.TFrame", padding=1)
            self.drop_wrap.pack(fill="both", expand=True, pady=(10, 0))
            self.drop_canvas = tk.Canvas(self.drop_wrap, bg=C["drop"],
                                         highlightthickness=0, bd=0, height=280,
                                         cursor="hand2")
            self.drop_canvas.pack(fill="both", expand=True)
            self.drop_canvas.bind("<Button-1>", lambda _e: self.conv_add_files())
            self.drop_canvas.bind("<Configure>", lambda _e: self._draw_dropzone())
            if drag_and_drop:
                self.drop_canvas.drop_target_register(DND_FILES)
                self.drop_canvas.dnd_bind("<<Drop>>", self._on_conv_drop)

            dest = ttk.Frame(parent, style="Card.TFrame", padding=12)
            dest.pack(fill="x", pady=(12, 0))
            ttk.Label(dest, text="Save to:", style="CardInk.TLabel").pack(side="left")
            self.conv_out_var = tk.StringVar(
                value=getattr(self, "_conv_out_saved", ""))
            entry = ttk.Entry(dest, textvariable=self.conv_out_var)
            entry.pack(side="left", fill="x", expand=True, padx=(10, 8))
            ttk.Button(dest, text="Browse…", style="Ghost.TButton",
                       command=self.conv_choose_folder).pack(side="left")
            self.convert_button = ttk.Button(dest, text="Convert Now",
                                              style="Primary.TButton",
                                              command=self.start_convert)
            self.convert_button.pack(side="left", padx=(10, 0))

            self.conv_status = tk.StringVar(value="No files added yet.")
            ttk.Label(parent, textvariable=self.conv_status, style="Status.TLabel"
                      ).pack(anchor="w", pady=(10, 0))

            self.conv_progress = ttk.Progressbar(
                parent, mode="determinate",
                style="Accent.Horizontal.TProgressbar")
            self._conv_route_changed(initial=True)

        def _draw_dropzone(self):
            if not hasattr(self, "drop_canvas") or not self.drop_canvas.winfo_exists():
                return
            C = self.C
            cv = self.drop_canvas
            cv.delete("all")
            w = cv.winfo_width() or cv.winfo_reqwidth()
            h = cv.winfo_height() or 280
            cv.create_rectangle(14, 14, w - 14, h - 14, outline=C["drop_bd"],
                                dash=(6, 4), width=2)
            if self.conv_inputs:
                cx = w // 2
                y = 44
                cv.create_text(cx, y, text=f"{len(self.conv_inputs)} file"
                               f"{'s' if len(self.conv_inputs) != 1 else ''} ready",
                               fill=C["ink"], font=(UI, 12, "bold"))
                y += 26
                for path in self.conv_inputs[:8]:
                    cv.create_text(cx, y, text=ellipsize(os.path.basename(path), 60),
                                   fill=C["muted"], font=(UI, 10))
                    y += 20
                if len(self.conv_inputs) > 8:
                    cv.create_text(cx, y, text=f"+ {len(self.conv_inputs) - 8} more",
                                   fill=C["muted"], font=(UI, 9))
                    y += 20
                cv.create_text(cx, h - 34, text="Click to add more · "
                               "Clear removes all", fill=C["muted"], font=(UI, 9))
                cv.bind("<Double-Button-1>", lambda _e: self.conv_clear())
                return
            cx, cy = w // 2, h // 2 - 20
            cv.create_oval(cx - 46, cy - 6, cx - 6, cy + 26, outline=C["accent"],
                           width=2)
            cv.create_oval(cx - 20, cy - 22, cx + 24, cy + 20, outline=C["accent"],
                           width=2)
            cv.create_oval(cx + 6, cy - 6, cx + 46, cy + 26, outline=C["accent"],
                           width=2)
            cv.create_rectangle(cx - 30, cy + 16, cx + 30, cy + 30, fill=C["drop"],
                                outline=C["drop"])
            cv.create_line(cx, cy + 6, cx, cy + 34, fill=C["accent"], width=2)
            cv.create_line(cx - 8, cy + 14, cx, cy + 6, fill=C["accent"], width=2)
            cv.create_line(cx + 8, cy + 14, cx, cy + 6, fill=C["accent"], width=2)
            cv.create_text(cx, cy + 70,
                           text="Drag and drop files here to convert, or click to "
                           "browse.", fill=C["muted"], font=(UI, 11))

        def _conv_route_changed(self, initial=False):
            route = self._conv_route()
            if route.available:
                exts = " ".join("*" + e for e in route.src_exts)
                self.conv_hint.set(f"Accepts {exts}. Output: {route.dst_ext}")
            else:
                self.conv_hint.set("Unavailable – " + route.hint)
            if not initial:
                self.conv_inputs.clear()
            self.conv_refresh()

        def conv_add_files(self):
            route = self._conv_route()
            if not route.available:
                messagebox.showinfo("Not available", route.hint)
                return
            patterns = ";".join("*" + e for e in route.src_exts)
            paths = filedialog.askopenfilenames(
                title="Add files to convert",
                filetypes=[(route.label, patterns), ("All files", "*.*")])
            self._conv_add_paths(paths)

        def _on_conv_drop(self, event):
            self._conv_add_paths(self.drop_canvas.tk.splitlist(event.data))

        def _conv_add_paths(self, paths):
            route = self._conv_route()
            skipped = []
            for path in paths:
                path = os.path.abspath(path)
                if (not os.path.isfile(path) or
                        os.path.splitext(path)[1].lower() not in route.src_exts):
                    skipped.append(os.path.basename(path))
                    continue
                if path not in self.conv_inputs:
                    self.conv_inputs.append(path)
            self.conv_refresh()
            if skipped:
                messagebox.showwarning(
                    "Some files were skipped",
                    f"These do not match {route.label}:\n\n"
                    + "\n".join(skipped[:10]))

        def conv_clear(self):
            self.conv_inputs.clear()
            self.conv_refresh()

        def conv_choose_folder(self):
            folder = filedialog.askdirectory(title="Save converted files to")
            if folder:
                self.conv_out_var.set(folder)

        def conv_refresh(self):
            self._draw_dropzone()
            route = self._conv_route()
            can = route.available and bool(self.conv_inputs) and not self.conv_working
            self.convert_button.state(["!disabled"] if can else ["disabled"])
            if not route.available:
                self.conv_status.set(f"{route.label} is unavailable on this machine.")
            elif not self.conv_inputs:
                self.conv_status.set("No files added yet.")
            else:
                n = len(self.conv_inputs)
                self.conv_status.set(f"{n} file{'s' if n != 1 else ''} ready.")

        def start_convert(self):
            if self.conv_working:
                return
            route = self._conv_route()
            if not route.available or not self.conv_inputs:
                return
            out_dir = self.conv_out_var.get().strip()
            self._conv_out_saved = out_dir
            if out_dir and not os.path.isdir(out_dir):
                if messagebox.askyesno("Create folder?",
                                       f"{out_dir}\n\ndoes not exist. Create it?"):
                    try:
                        os.makedirs(out_dir, exist_ok=True)
                    except OSError as exc:
                        messagebox.showerror("Could not create folder", str(exc))
                        return
                else:
                    return
            targets = [route.output_for(s, out_dir) for s in self.conv_inputs]
            existing = [os.path.basename(t) for t in targets if os.path.exists(t)]
            if existing and not messagebox.askyesno(
                    "Replace files?", "These already exist and will be replaced:\n\n"
                    + "\n".join(existing[:10]) + "\n\nContinue?"):
                return

            self.conv_working = True
            self.convert_button.state(["disabled"])
            self.conv_progress.pack(fill="x", pady=(10, 0))
            self.conv_progress.configure(value=0, maximum=len(self.conv_inputs))
            self.conv_status.set("Converting…")
            inputs = list(self.conv_inputs)
            route_key = route.key

            def work():
                if HAS_COMTYPES:
                    try:
                        import comtypes
                        comtypes.CoInitialize()
                    except Exception:
                        pass
                the_route = route_by_key(route_key)
                ok = 0
                try:
                    for number, src in enumerate(inputs, start=1):
                        base = os.path.basename(src)
                        self.conv_events.put(("file", number, len(inputs), base))

                        def progress(done, total, name, _n=number, _b=base):
                            self.conv_events.put(("page", _n, len(inputs),
                                                  f"{_b} - {name}" if name else _b,
                                                  done, total))
                        try:
                            written = convert_file(src, the_route, out_dir=out_dir,
                                                   progress=progress)
                            self.conv_events.put(("ok", number, src, written))
                            ok += 1
                        except ConversionError as exc:
                            self.conv_events.put(("fail", number, base, str(exc)))
                        except Exception as exc:
                            self.conv_events.put(("fail", number, base,
                                                  f"Unexpected: {exc}"))
                finally:
                    if HAS_COMTYPES:
                        try:
                            import comtypes
                            comtypes.CoUninitialize()
                        except Exception:
                            pass
                    self.conv_events.put(("finished", ok, len(inputs), out_dir))

            threading.Thread(target=work, daemon=True).start()

        def _drain_conv_events(self):
            try:
                while True:
                    event = self.conv_events.get_nowait()
                    kind = event[0]
                    if kind == "file":
                        _, number, total, base = event
                        self.conv_progress.configure(maximum=total, value=number - 1)
                        self.conv_status.set(f"[{number}/{total}] Converting {base}…")
                    elif kind == "page":
                        _, number, total, label, done, page_total = event
                        detail = (f" · page {done}/{page_total}"
                                  if page_total and page_total > 1 else "")
                        self.conv_status.set(f"[{number}/{total}] {label}{detail}")
                    elif kind == "ok":
                        _, number, _s, _w = event
                        self.conv_progress.configure(value=number)
                    elif kind == "fail":
                        _, number, base, error = event
                        self.conv_progress.configure(value=number)
                        messagebox.showerror("Conversion failed", f"{base}:\n\n{error}")
                    else:
                        _, ok, total, out_dir = event
                        self.conv_working = False
                        failed = total - ok
                        summary = (f"Converted {ok} of {total} file"
                                   f"{'s' if total != 1 else ''}.")
                        if failed:
                            summary += f" {failed} failed."
                        self.conv_status.set(summary)
                        self.conv_refresh()
                        if ok:
                            where = out_dir or (os.path.dirname(self.conv_inputs[0])
                                                if self.conv_inputs else "")
                            if where and messagebox.askyesno(
                                    "Done", summary + "\n\nOpen the output folder?"):
                                open_folder(os.path.join(where, "_"))
            except queue.Empty:
                pass
            self._conv_drain_id = self.after(100, self._drain_conv_events)

        # ==============================================================
        # EXTRACT TAB
        # ==============================================================
        def _build_extract(self, parent):
            C = self.C
            card = ttk.Frame(parent, style="Card.TFrame", padding=16)
            card.pack(fill="x")
            card.columnconfigure(1, weight=1)

            ttk.Label(card, text="Source PDF:", style="CardInk.TLabel").grid(
                row=0, column=0, sticky="w", pady=(0, 10))
            self.ex_src = tk.StringVar(value=getattr(self, "_ex_src_saved", ""))
            ttk.Entry(card, textvariable=self.ex_src).grid(
                row=0, column=1, sticky="ew", padx=(10, 8), pady=(0, 10))
            ttk.Button(card, text="Browse…", style="Ghost.TButton",
                       command=self._ex_pick_source).grid(row=0, column=2,
                                                          pady=(0, 10))

            ttk.Label(card, text="Pages:", style="CardInk.TLabel").grid(
                row=1, column=0, sticky="w", pady=(0, 10))
            self.ex_pages = tk.StringVar(value=getattr(self, "_ex_pages_saved", "all"))
            pages_entry = ttk.Entry(card, textvariable=self.ex_pages)
            pages_entry.grid(row=1, column=1, sticky="w", padx=(10, 8), pady=(0, 10))
            pages_entry.configure(width=24)
            ttk.Label(card, text="e.g. 1-3, 7  ·  all", style="CardMuted.TLabel"
                      ).grid(row=1, column=1, sticky="w", padx=(210, 0), pady=(0, 10))

            ttk.Label(card, text="Extract:", style="CardInk.TLabel").grid(
                row=2, column=0, sticky="w")
            seg = ttk.Frame(card, style="Card.TFrame")
            seg.grid(row=2, column=1, sticky="w", padx=(10, 0))
            self.ex_mode = getattr(self, "ex_mode", "pages")
            self._ex_seg_buttons = {}
            for key, text in (("pages", "Pages → PDF"),
                              ("text", "Text → .txt"),
                              ("images", "Images → folder")):
                btn = ttk.Button(seg, text=text,
                                 style="SegOn.TButton" if self.ex_mode == key
                                 else "Seg.TButton",
                                 command=lambda k=key: self._ex_set_mode(k))
                btn.pack(side="left", padx=(0, 8))
                self._ex_seg_buttons[key] = btn

            dest = ttk.Frame(parent, style="Card.TFrame", padding=12)
            dest.pack(fill="x", pady=(12, 0))
            ttk.Label(dest, text="Save to:", style="CardInk.TLabel").pack(side="left")
            self.ex_out = tk.StringVar(value=getattr(self, "_ex_out_saved", ""))
            ttk.Entry(dest, textvariable=self.ex_out).pack(
                side="left", fill="x", expand=True, padx=(10, 8))
            ttk.Button(dest, text="Browse…", style="Ghost.TButton",
                       command=self._ex_pick_folder).pack(side="left")
            self.extract_button = ttk.Button(dest, text="Extract",
                                              style="Primary.TButton",
                                              command=self.start_extract)
            self.extract_button.pack(side="left", padx=(10, 0))

            ttk.Label(parent, text="Leave 'Save to' blank to write next to the "
                      "source PDF.", style="Muted.TLabel").pack(anchor="w",
                                                                pady=(6, 0))
            self.ex_progress = ttk.Progressbar(
                parent, mode="determinate",
                style="Accent.Horizontal.TProgressbar")
            self.ex_status = tk.StringVar(value="Choose a PDF to extract from.")
            ttk.Label(parent, textvariable=self.ex_status, style="Status.TLabel"
                      ).pack(anchor="w", pady=(10, 0))

        def _ex_set_mode(self, key):
            self.ex_mode = key
            for name, btn in self._ex_seg_buttons.items():
                btn.configure(style="SegOn.TButton" if name == key
                              else "Seg.TButton")

        def _ex_pick_source(self):
            path = filedialog.askopenfilename(
                title="Choose a PDF", filetypes=[("PDF files", "*.pdf")])
            if path:
                self.ex_src.set(path)
                self.ex_status.set(f"Ready: {os.path.basename(path)}")

        def _ex_pick_folder(self):
            folder = filedialog.askdirectory(title="Save extracted output to")
            if folder:
                self.ex_out.set(folder)

        def start_extract(self):
            if self.ex_working:
                return
            src = self.ex_src.get().strip()
            if not src or not os.path.isfile(src):
                messagebox.showinfo("Choose a PDF", "Pick a source PDF first.")
                return
            self._ex_src_saved = src
            self._ex_pages_saved = self.ex_pages.get().strip() or "all"
            self._ex_out_saved = self.ex_out.get().strip()
            pages = self._ex_pages_saved
            out_dir = self._ex_out_saved or os.path.dirname(os.path.abspath(src))
            stem = os.path.splitext(os.path.basename(src))[0]
            mode = self.ex_mode

            if mode == "pages":
                target = os.path.join(out_dir, f"{stem}_extracted.pdf")
                job = lambda p: ("pages", extract_pages(src, target, pages,
                                                        progress=p), target)
            elif mode == "text":
                target = os.path.join(out_dir, f"{stem}.txt")
                job = lambda p: ("text", extract_text(src, target, pages,
                                                      progress=p), target)
            else:
                target = os.path.join(out_dir, f"{stem}_images")
                job = lambda p: ("images", extract_images(src, target, pages,
                                                         progress=p), target)

            self.ex_working = True
            self.extract_button.state(["disabled"])
            self.ex_progress.pack(fill="x", pady=(10, 0))
            self.ex_progress.configure(value=0, maximum=100)
            self.ex_status.set("Extracting…")

            def work():
                try:
                    kind, amount, out = job(
                        lambda done, total, name:
                        self.ex_events.put(("progress", done, total, name)))
                    self.ex_events.put(("done", kind, amount, out, None))
                except (ExtractError, MergeError) as exc:
                    self.ex_events.put(("done", "", 0, "", str(exc)))
                except Exception as exc:
                    self.ex_events.put(("done", "", 0, "",
                                        f"Unexpected problem: {exc}"))

            threading.Thread(target=work, daemon=True).start()

        def _drain_ex_events(self):
            try:
                while True:
                    event = self.ex_events.get_nowait()
                    if event[0] == "progress":
                        _, done, total, name = event
                        self.ex_progress.configure(maximum=max(total, 1), value=done)
                        self.ex_status.set(f"{name} · {done}/{total}")
                    else:
                        _, kind, amount, out, error = event
                        self.ex_working = False
                        self.extract_button.state(["!disabled"])
                        if error:
                            self.ex_progress.configure(value=0)
                            self.ex_status.set("Extract failed.")
                            messagebox.showerror("Extract failed", error)
                        else:
                            if kind == "pages":
                                msg = f"Extracted {amount} pages to {os.path.basename(out)}"
                            elif kind == "text":
                                msg = f"Wrote {amount} characters to {os.path.basename(out)}"
                            else:
                                msg = f"Saved {amount} image{'s' if amount != 1 else ''} to {os.path.basename(out)}\\"
                            self.ex_progress.configure(value=self.ex_progress["maximum"])
                            self.ex_status.set(msg)
                            if messagebox.askyesno("Extracted",
                                                   msg + "\n\nOpen the folder?"):
                                open_folder(out if os.path.isfile(out)
                                            else os.path.join(out, "_"))
            except queue.Empty:
                pass
            self._ex_drain_id = self.after(100, self._drain_ex_events)

        # ==============================================================
        # WATERMARK TAB
        # ==============================================================
        def _build_watermark(self, parent):
            C = self.C
            top = ttk.Frame(parent, style="App.TFrame")
            top.pack(fill="x")
            top.columnconfigure(0, weight=1)

            form = ttk.Frame(top, style="Card.TFrame", padding=16)
            form.grid(row=0, column=0, sticky="nsew")
            form.columnconfigure(1, weight=1)

            ttk.Label(form, text="Source:", style="CardInk.TLabel").grid(
                row=0, column=0, sticky="w", pady=(0, 10))
            self.wm_src = tk.StringVar(value=getattr(self, "_wm_src_saved", ""))
            ttk.Entry(form, textvariable=self.wm_src).grid(
                row=0, column=1, sticky="ew", padx=(10, 8), pady=(0, 10))
            ttk.Button(form, text="Browse…", style="Ghost.TButton",
                       command=self._wm_pick_source).grid(row=0, column=2,
                                                          pady=(0, 10))

            ttk.Label(form, text="Type:", style="CardInk.TLabel").grid(
                row=1, column=0, sticky="w", pady=(0, 10))
            typ = ttk.Frame(form, style="Card.TFrame")
            typ.grid(row=1, column=1, columnspan=2, sticky="w", padx=(10, 0),
                     pady=(0, 10))
            self.wm_type = getattr(self, "wm_type", "text")
            self._wm_typeseg = {}
            for key, label in (("text", "Text"), ("image", "Image / logo")):
                btn = ttk.Button(typ, text=label,
                                 style="SegOn.TButton" if self.wm_type == key
                                 else "Seg.TButton",
                                 command=lambda k=key: self._wm_set_type(k))
                btn.pack(side="left", padx=(0, 8))
                self._wm_typeseg[key] = btn

            # text row and image row share grid row 2 (only one shown at a time)
            self._wm_text_row = ttk.Frame(form, style="Card.TFrame")
            self._wm_text_row.grid(row=2, column=0, columnspan=3, sticky="ew",
                                   pady=(0, 10))
            self._wm_text_row.columnconfigure(1, weight=1)
            ttk.Label(self._wm_text_row, text="Watermark text:",
                      style="CardInk.TLabel").grid(row=0, column=0, sticky="w")
            self.wm_text = tk.StringVar(
                value=getattr(self, "_wm_text_saved", "CONFIDENTIAL"))
            te = ttk.Entry(self._wm_text_row, textvariable=self.wm_text)
            te.grid(row=0, column=1, sticky="ew", padx=(10, 0))
            te.bind("<KeyRelease>", lambda _e: self._wm_schedule_preview())

            self._wm_image_row = ttk.Frame(form, style="Card.TFrame")
            self._wm_image_row.grid(row=2, column=0, columnspan=3, sticky="ew",
                                    pady=(0, 10))
            self._wm_image_row.columnconfigure(1, weight=1)
            ttk.Label(self._wm_image_row, text="Image / logo:",
                      style="CardInk.TLabel").grid(row=0, column=0, sticky="w")
            self.wm_image = tk.StringVar(value=getattr(self, "_wm_image_saved", ""))
            ttk.Entry(self._wm_image_row, textvariable=self.wm_image).grid(
                row=0, column=1, sticky="ew", padx=(10, 8))
            ttk.Button(self._wm_image_row, text="Browse…", style="Ghost.TButton",
                       command=self._wm_pick_image).grid(row=0, column=2)

            self._wm_color_row = ttk.Frame(form, style="Card.TFrame")
            self._wm_color_row.grid(row=3, column=0, columnspan=3, sticky="w",
                                    pady=(0, 10))
            ttk.Label(self._wm_color_row, text="Colour:",
                      style="CardInk.TLabel").pack(side="left")
            self.wm_color = tk.StringVar(value=getattr(self, "_wm_color_saved",
                                                       "Gray"))
            cbox = ttk.Combobox(self._wm_color_row, textvariable=self.wm_color,
                                state="readonly", width=8,
                                values=list(WATERMARK_COLORS) + ["Custom"])
            cbox.pack(side="left", padx=(10, 8))
            cbox.bind("<<ComboboxSelected>>", lambda _e: self._wm_color_changed())
            self.wm_swatch = tk.Frame(self._wm_color_row, width=22, height=22,
                                      highlightthickness=1,
                                      highlightbackground=C["border"])
            self.wm_swatch.pack(side="left")
            self.wm_swatch.pack_propagate(False)
            ttk.Button(self._wm_color_row, text="Custom…", style="Ghost.TButton",
                       command=self._wm_pick_color).pack(side="left", padx=(8, 0))

            ttk.Label(form, text="Style:", style="CardInk.TLabel").grid(
                row=4, column=0, sticky="w", pady=(0, 10))
            strow = ttk.Frame(form, style="Card.TFrame")
            strow.grid(row=4, column=1, columnspan=2, sticky="w", padx=(10, 0),
                       pady=(0, 10))
            self.wm_opacity = tk.StringVar(value=getattr(self, "_wm_opacity_saved",
                                                         "15%"))
            ttk.Label(strow, text="Opacity", style="CardMuted.TLabel").pack(
                side="left")
            ocb = ttk.Combobox(strow, textvariable=self.wm_opacity,
                               state="readonly", width=6,
                               values=["10%", "15%", "20%", "30%", "50%", "70%",
                                       "100%"])
            ocb.pack(side="left", padx=(6, 16))
            ocb.bind("<<ComboboxSelected>>", lambda _e: self._wm_schedule_preview())
            self.wm_rot = tk.StringVar(value=getattr(self, "_wm_rot_saved", "45°"))
            ttk.Label(strow, text="Angle", style="CardMuted.TLabel").pack(
                side="left")
            acb = ttk.Combobox(strow, textvariable=self.wm_rot, state="readonly",
                               width=6, values=["45°", "30°", "0°", "90°", "-45°"])
            acb.pack(side="left", padx=(6, 16))
            acb.bind("<<ComboboxSelected>>", lambda _e: self._wm_schedule_preview())
            self.wm_size = tk.StringVar(value=getattr(self, "_wm_size_saved", "48"))
            self._wm_size_lbl = ttk.Label(strow, text="Size",
                                          style="CardMuted.TLabel")
            self._wm_size_lbl.pack(side="left")
            sze = ttk.Entry(strow, textvariable=self.wm_size, width=6)
            sze.pack(side="left", padx=(6, 0))
            sze.bind("<KeyRelease>", lambda _e: self._wm_schedule_preview())

            ttk.Label(form, text="Layout:", style="CardInk.TLabel").grid(
                row=5, column=0, sticky="w")
            lay = ttk.Frame(form, style="Card.TFrame")
            lay.grid(row=5, column=1, columnspan=2, sticky="w", padx=(10, 0))
            self.wm_layout = getattr(self, "wm_layout", "centered")
            self._wm_seg = {}
            for key, label in (("centered", "Centered"), ("tiled", "Tiled")):
                btn = ttk.Button(lay, text=label,
                                 style="SegOn.TButton" if self.wm_layout == key
                                 else "Seg.TButton",
                                 command=lambda k=key: self._wm_set_layout(k))
                btn.pack(side="left", padx=(0, 8))
                self._wm_seg[key] = btn
            ttk.Label(lay, text="    Pages:", style="CardMuted.TLabel").pack(
                side="left", padx=(10, 0))
            self.wm_pages = tk.StringVar(value=getattr(self, "_wm_pages_saved",
                                                       "all"))
            pen = ttk.Entry(lay, textvariable=self.wm_pages, width=12)
            pen.pack(side="left", padx=(6, 0))
            pen.bind("<KeyRelease>", lambda _e: self._wm_schedule_preview())
            self._wm_pages_entry = pen

            # tiled-only: how many copies across the page
            self._wm_density_lbl = ttk.Label(lay, text="    Tiles across:",
                                             style="CardMuted.TLabel")
            self.wm_density = tk.StringVar(value=getattr(self, "_wm_density_saved",
                                                         "4"))
            self._wm_density_spin = ttk.Spinbox(
                lay, from_=1, to=16, width=4, textvariable=self.wm_density,
                command=self._wm_schedule_preview)
            self._wm_density_spin.bind("<KeyRelease>",
                                       lambda _e: self._wm_schedule_preview())

            # live preview (right column)
            prev = ttk.Frame(top, style="Card.TFrame", padding=12)
            prev.grid(row=0, column=1, sticky="nsew", padx=(12, 0))
            ttk.Label(prev, text="Live preview", style="CardInk.TLabel").pack(
                anchor="w")
            self.wm_preview_holder = tk.Frame(prev, bg=C["page"], width=300,
                                              height=400, highlightthickness=1,
                                              highlightbackground=C["border"])
            self.wm_preview_holder.pack(pady=(8, 0))
            self.wm_preview_holder.pack_propagate(False)
            self.wm_preview_label = tk.Label(self.wm_preview_holder, bg=C["page"],
                                             fg=C["muted"], font=(UI, 9),
                                             text="Select a PDF to preview")
            self.wm_preview_label.pack(fill="both", expand=True)
            self.wm_preview_cap = tk.StringVar(value="")
            ttk.Label(prev, textvariable=self.wm_preview_cap,
                      style="CardMuted.TLabel").pack(anchor="w", pady=(6, 0))

            dest = ttk.Frame(parent, style="Card.TFrame", padding=12)
            dest.pack(fill="x", pady=(12, 0))
            ttk.Label(dest, text="Save to:", style="CardInk.TLabel").pack(side="left")
            self.wm_out = tk.StringVar(value=getattr(self, "_wm_out_saved", ""))
            ttk.Entry(dest, textvariable=self.wm_out).pack(
                side="left", fill="x", expand=True, padx=(10, 8))
            ttk.Button(dest, text="Browse…", style="Ghost.TButton",
                       command=self._wm_pick_folder).pack(side="left")
            self.wm_button = ttk.Button(dest, text="Apply Watermark",
                                        style="Primary.TButton",
                                        command=self.start_watermark)
            self.wm_button.pack(side="left", padx=(10, 0))

            ttk.Label(parent, text="Source can be a PDF or an image. Leave "
                      "'Save to' blank to write the result next to the source.",
                      style="Muted.TLabel").pack(anchor="w", pady=(6, 0))
            self.wm_progress = ttk.Progressbar(
                parent, mode="determinate",
                style="Accent.Horizontal.TProgressbar")
            self.wm_status = tk.StringVar(value="Choose a PDF or image to "
                                          "watermark.")
            ttk.Label(parent, textvariable=self.wm_status, style="Status.TLabel"
                      ).pack(anchor="w", pady=(10, 0))

            self._wm_custom = getattr(self, "_wm_custom", None)
            self._wm_preview_after = None
            self._wm_apply_type_visibility()
            self._wm_update_swatch()
            self._wm_update_source_kind()
            self._wm_update_density_vis()
            self._wm_schedule_preview()

        def _wm_update_density_vis(self):
            if self.wm_layout == "tiled":
                self._wm_density_lbl.pack(side="left", padx=(10, 0))
                self._wm_density_spin.pack(side="left", padx=(6, 0))
            else:
                self._wm_density_lbl.pack_forget()
                self._wm_density_spin.pack_forget()

        def _wm_update_source_kind(self):
            src = self.wm_src.get().strip()
            is_img = bool(src) and is_watermark_image(src)
            try:
                self._wm_pages_entry.configure(
                    state="disabled" if is_img else "normal")
            except Exception:
                pass

        def _wm_apply_type_visibility(self):
            if self.wm_type == "text":
                self._wm_image_row.grid_remove()
                self._wm_text_row.grid()
                self._wm_color_row.grid()
                self._wm_size_lbl.configure(text="Size")
            else:
                self._wm_text_row.grid_remove()
                self._wm_image_row.grid()
                self._wm_color_row.grid_remove()
                self._wm_size_lbl.configure(text="Width")

        def _wm_set_type(self, key):
            if key == self.wm_type:
                return
            current = self.wm_size.get().strip()
            if key == "image" and current in ("", "48"):
                self.wm_size.set("220")
            if key == "text" and current in ("", "220"):
                self.wm_size.set("48")
            self.wm_type = key
            for name, btn in self._wm_typeseg.items():
                btn.configure(style="SegOn.TButton" if name == key
                              else "Seg.TButton")
            self._wm_apply_type_visibility()
            self._wm_schedule_preview()

        def _wm_set_layout(self, key):
            self.wm_layout = key
            for name, btn in self._wm_seg.items():
                btn.configure(style="SegOn.TButton" if name == key
                              else "Seg.TButton")
            self._wm_update_density_vis()
            self._wm_schedule_preview()

        def _wm_pick_source(self):
            images = "*.png *.jpg *.jpeg *.bmp *.gif *.tif *.tiff *.webp"
            path = filedialog.askopenfilename(
                title="Choose a PDF or image",
                filetypes=[("PDF or image", "*.pdf " + images),
                           ("PDF files", "*.pdf"), ("Images", images),
                           ("All files", "*.*")])
            if path:
                self.wm_src.set(path)
                self.wm_status.set(f"Ready: {os.path.basename(path)}")
                self._wm_update_source_kind()
                self._wm_schedule_preview()

        def _wm_pick_image(self):
            path = filedialog.askopenfilename(
                title="Choose a watermark image",
                filetypes=[("Images", "*.png *.jpg *.jpeg *.gif *.bmp"),
                           ("All files", "*.*")])
            if path:
                self.wm_image.set(path)
                self._wm_schedule_preview()

        def _wm_pick_folder(self):
            folder = filedialog.askdirectory(title="Save watermarked PDF to")
            if folder:
                self.wm_out.set(folder)

        def _wm_active_color(self):
            if self.wm_color.get() == "Custom" and self._wm_custom:
                return self._wm_custom
            return WATERMARK_COLORS.get(self.wm_color.get(), (0.5, 0.5, 0.5))

        def _wm_update_swatch(self):
            r, g, b = self._wm_active_color()
            self.wm_swatch.configure(
                bg="#%02x%02x%02x" % (int(r * 255), int(g * 255), int(b * 255)))

        def _wm_color_changed(self):
            if self.wm_color.get() == "Custom" and not self._wm_custom:
                self._wm_pick_color()
            else:
                self._wm_update_swatch()
                self._wm_schedule_preview()

        def _wm_pick_color(self):
            init = self._wm_active_color()
            init_hex = "#%02x%02x%02x" % (int(init[0] * 255), int(init[1] * 255),
                                          int(init[2] * 255))
            rgb, _hex = colorchooser.askcolor(color=init_hex,
                                              title="Watermark colour", parent=self)
            if rgb:
                self._wm_custom = (round(rgb[0]) / 255, round(rgb[1]) / 255,
                                   round(rgb[2]) / 255)
                self.wm_color.set("Custom")
                self._wm_update_swatch()
                self._wm_schedule_preview()

        def _wm_gather(self):
            try:
                size = float(self.wm_size.get().strip() or "48")
            except ValueError:
                size = 48.0
            opacity = int("".join(c for c in self.wm_opacity.get()
                                  if c.isdigit()) or "15") / 100.0
            opacity = max(0.0, min(1.0, opacity))
            try:
                angle = float(self.wm_rot.get().replace("°", "") or "45")
            except ValueError:
                angle = 45.0
            try:
                density = max(1, min(20, int(float(self.wm_density.get() or "4"))))
            except (ValueError, AttributeError):
                density = 4
            return dict(size=size, opacity=opacity, angle=angle, density=density,
                        color=self._wm_active_color(),
                        tiled=(self.wm_layout == "tiled"), type=self.wm_type,
                        text=self.wm_text.get(), image=self.wm_image.get().strip(),
                        pages=self.wm_pages.get().strip() or "all")

        def _wm_schedule_preview(self):
            if getattr(self, "_wm_preview_after", None):
                try:
                    self.after_cancel(self._wm_preview_after)
                except Exception:
                    pass
            self._wm_preview_after = self.after(250, self._wm_render_preview)

        def _wm_render_preview(self):
            self._wm_preview_after = None
            if not hasattr(self, "wm_preview_label"):
                return
            try:
                if not self.wm_preview_label.winfo_exists():
                    return
            except Exception:
                return
            if not HAS_FITZ:
                self.wm_preview_label.configure(image="", text="Preview needs "
                                                "PyMuPDF")
                return
            src = self.wm_src.get().strip()
            if not src or not os.path.isfile(src):
                self.wm_preview_label.configure(
                    image="", text="Select a PDF or image to preview")
                self.wm_preview_cap.set("")
                return
            settings = self._wm_gather()
            try:
                if is_watermark_image(src):
                    single, page, _scale = _image_to_page(src)
                    source_label = "Image"
                else:
                    document = _fitz.open(src)
                    total = document.page_count
                    try:
                        indices = parse_page_range(settings["pages"], total)
                    except MergeError:
                        indices = list(range(total))
                    page_index = indices[0] if indices else 0
                    single = _fitz.open()
                    single.insert_pdf(document, from_page=page_index,
                                      to_page=page_index)
                    document.close()
                    page = single[0]
                    source_label = f"Page {page_index + 1}"
                if (settings["type"] == "image" and settings["image"]
                        and os.path.isfile(settings["image"])):
                    png, (iw, ih) = _prepare_wm_image(
                        settings["image"], settings["opacity"], settings["angle"])
                    _stamp_image_page(page, png, iw, ih, settings["size"],
                                      settings["tiled"], settings["density"])
                    kind = "image"
                elif settings["type"] == "text" and settings["text"].strip():
                    _stamp_page(page, settings["text"], settings["size"],
                                settings["color"], settings["opacity"],
                                settings["angle"], settings["tiled"],
                                settings["density"])
                    kind = "text"
                else:
                    kind = "none"
                zoom = min(286 / max(page.rect.width, 1),
                           384 / max(page.rect.height, 1))
                pixmap = page.get_pixmap(matrix=_fitz.Matrix(zoom, zoom),
                                         alpha=False)
                single.close()
                photo = tk.PhotoImage(
                    master=self,
                    data=base64.b64encode(pixmap.tobytes("png")).decode("ascii"))
                self._wm_preview_img = photo
                self.wm_preview_label.configure(image=photo, text="")
                label = {"image": "image watermark", "text": "text watermark",
                         "none": "no watermark yet"}[kind]
                self.wm_preview_cap.set(f"{source_label} · {label}")
            except Exception:
                try:
                    self.wm_preview_label.configure(image="",
                                                    text="Preview unavailable")
                except Exception:
                    pass

        def start_watermark(self):
            if self.wm_working:
                return
            src = self.wm_src.get().strip()
            if not src or not os.path.isfile(src):
                messagebox.showinfo("Choose a file",
                                    "Pick a source PDF or image first.")
                return
            settings = self._wm_gather()
            if settings["type"] == "image":
                if not settings["image"] or not os.path.isfile(settings["image"]):
                    messagebox.showinfo("Choose an image",
                                        "Pick a watermark image (PNG/JPG).")
                    return
            elif not settings["text"].strip():
                messagebox.showinfo("Add text", "Enter the watermark text.")
                return

            out_dir = self.wm_out.get().strip() or os.path.dirname(
                os.path.abspath(src))
            stem = os.path.splitext(os.path.basename(src))[0]
            out_ext = (os.path.splitext(src)[1].lower()
                       if is_watermark_image(src) else ".pdf")
            target = os.path.join(out_dir, f"{stem}_watermarked{out_ext}")

            self._wm_src_saved = src
            self._wm_text_saved = self.wm_text.get()
            self._wm_image_saved = settings["image"]
            self._wm_color_saved = self.wm_color.get()
            self._wm_opacity_saved = self.wm_opacity.get()
            self._wm_rot_saved = self.wm_rot.get()
            self._wm_size_saved = self.wm_size.get()
            self._wm_pages_saved = settings["pages"]
            self._wm_density_saved = str(settings["density"])
            self._wm_out_saved = self.wm_out.get().strip()

            self.wm_working = True
            self.wm_button.state(["disabled"])
            self.wm_progress.pack(fill="x", pady=(10, 0))
            self.wm_progress.configure(value=0, maximum=100)
            self.wm_status.set("Applying watermark…")

            text = settings["text"]
            image = settings["image"] if settings["type"] == "image" else ""
            size = settings["size"]
            color = settings["color"]
            opacity = settings["opacity"]
            angle = settings["angle"]
            tiled = settings["tiled"]
            pages = settings["pages"]
            density = settings["density"]

            def work():
                try:
                    stamped = apply_watermark(
                        src, target, text=text, pages=pages, font_size=size,
                        color=color, opacity=opacity, angle=angle, tiled=tiled,
                        image=image, image_width=size, density=density,
                        progress=lambda d, t, nm:
                            self.wm_events.put(("progress", d, t, nm)))
                    self.wm_events.put(("done", stamped, target, None))
                except (WatermarkError, MergeError) as exc:
                    self.wm_events.put(("done", 0, "", str(exc)))
                except Exception as exc:
                    self.wm_events.put(("done", 0, "", f"Unexpected problem: {exc}"))

            threading.Thread(target=work, daemon=True).start()

        def _drain_wm_events(self):
            try:
                while True:
                    event = self.wm_events.get_nowait()
                    if event[0] == "progress":
                        _, done, total, name = event
                        self.wm_progress.configure(maximum=max(total, 1), value=done)
                        self.wm_status.set(f"{name} · {done}/{total}")
                    else:
                        _, stamped, out, error = event
                        self.wm_working = False
                        self.wm_button.state(["!disabled"])
                        if error:
                            self.wm_progress.configure(value=0)
                            self.wm_status.set("Watermark failed.")
                            messagebox.showerror("Watermark failed", error)
                        else:
                            self.wm_progress.configure(
                                value=self.wm_progress["maximum"])
                            if is_watermark_image(out):
                                what = "image"
                            else:
                                what = (f"{stamped} page"
                                        f"{'s' if stamped != 1 else ''}")
                            self.wm_status.set(
                                f"Watermarked {what} -> {os.path.basename(out)}")
                            if messagebox.askyesno(
                                    "Done", f"Watermarked {what} to:\n"
                                    f"{out}\n\nOpen the folder?"):
                                open_folder(out)
            except queue.Empty:
                pass
            self._wm_drain_id = self.after(100, self._drain_wm_events)

        # -- teardown --------------------------------------------------
        def destroy(self):
            for attr in ("_drain_id", "_conv_drain_id", "_ex_drain_id",
                         "_wm_drain_id", "_wm_preview_after"):
                ident = getattr(self, attr, None)
                if ident is not None:
                    try:
                        self.after_cancel(ident)
                    except Exception:
                        pass
                    setattr(self, attr, None)
            super().destroy()

    app = MantoolsApp()
    if _smoke is not None:
        try:
            _smoke(app)
        finally:
            try:
                app.destroy()
            except Exception:
                pass
        return 0
    app.mainloop()
    return 0


def open_folder(path: str) -> None:
    """Reveal a file in the system file manager. Never raises."""
    folder = os.path.dirname(os.path.abspath(path))
    try:
        if sys.platform.startswith("win"):
            os.startfile(folder)  # type: ignore[attr-defined]
        elif sys.platform == "darwin":
            import subprocess
            subprocess.Popen(["open", folder])
        else:
            import subprocess
            subprocess.Popen(["xdg-open", folder])
    except Exception:
        pass


def main() -> int:
    argv = sys.argv[1:]
    if argv and argv[0] == "convert":
        return run_convert_cli(argv[1:])
    if argv and argv[0] == "extract":
        return run_extract_cli(argv[1:])
    if argv and argv[0] == "watermark":
        return run_watermark_cli(argv[1:])
    if argv and argv[0] == "merge":
        return run_cli(argv[1:])
    if argv and argv[0] == "--cli":
        return run_cli(argv[1:])
    if argv:
        return run_cli(argv)
    return run_gui()


if __name__ == "__main__":
    sys.exit(main())
